import os
import time
import threading
import traceback
import json
import re
import hashlib
import urllib.request
import urllib.error
import sqlite3
from collections import Counter
from threading import Thread
from pathlib import Path
import logging

def _get_cv2():
    try:
        import cv2
        return cv2
    except ImportError:
        return None

# State Management
import backend.app.state as app_state
from backend.app.state import STATE
import backend.app.shared_state as shared_state
from backend.app.utils.log_utils import log_operation
from backend.app.utils.utils import parse_tags, _resolve_path

# Database & Config
from backend.app.database import SessionLocal, FileIndex
from backend.app.config import load_config
from sqlalchemy import func, text

try:
    from PIL import Image, ExifTags
except ImportError:
    Image = None
    ExifTags = None

def _get_fitz():
    with shared_state.MEMORY_LOCK:
        try:
            import fitz
            return fitz
        except ImportError:
            return None

def _get_mutagen():
    with shared_state.MEMORY_LOCK:
        try:
            import mutagen
            return mutagen
        except ImportError:
            return None

def _get_pefile():
    with shared_state.MEMORY_LOCK:
        try:
            import pefile
            return pefile
        except ImportError:
            return None

def _get_filetype():
    with shared_state.MEMORY_LOCK:
        try:
            import filetype
            return filetype
        except ImportError:
            return None

def _get_openpyxl():
    with shared_state.MEMORY_LOCK:
        try:
            import openpyxl
            return openpyxl
        except ImportError:
            return None

def _get_docx():
    with shared_state.MEMORY_LOCK:
        try:
            import docx
            return docx
        except ImportError:
            return None

def _get_pptx():
    with shared_state.MEMORY_LOCK:
        try:
            import pptx
            return pptx
        except ImportError:
            return None

STOP_WORDS = frozenset({
    "the", "and", "to", "of", "a", "in", "is", "that", "for", "it", "with", "as", "was", "on", 
    "at", "by", "an", "be", "this", "which", "or", "from", "but", "not", "are", "have", "we", 
    "all", "they", "one", "has", "were", "will", "would", "can", "if", "their", "your", "what", 
    "about", "who", "so", "when", "there", "out", "more", "do", "up", "any", "some", "into", 
    "only", "than", "them", "its", "also", "then", "been", "these", "how", "other", "could", 
    "our", "such", "may", "no", "over", "like", "most", "even", "because", "well", "where", 
    "through", "back", "much", "down", "should", "those", "under", "while", "very", "just", 
    "before", "each", "does", "why", "same", "both", "many", "after", "between", "too", "now", 
    "my", "me", "am", "i", "he", "she", "his", "hers", "him", "her", "us", "you", "yours"
})

IMAGENET_MAPPING = None

def load_imagenet_mapping():
    global IMAGENET_MAPPING
    if IMAGENET_MAPPING is None:
        try:
            # Try compiled python dict for maximum load speed
            from backend.app.utils.imagenet_mapping_data import IMAGENET_TO_COMMON
            IMAGENET_MAPPING = IMAGENET_TO_COMMON
        except Exception as e:
            print(f"Error loading ImageNet mapping: {e}")
            IMAGENET_MAPPING = {}
    return IMAGENET_MAPPING

def unload_imagenet_mapping():
    global IMAGENET_MAPPING
    IMAGENET_MAPPING = None
    import sys
    # De-register the compiled mapping module from sys.modules to release memory
    for target in ["backend.app.utils.imagenet_mapping_data", "backend.app.utils.imagenet_mapping_data."]:
        for mod in list(sys.modules.keys()):
            if mod == target or mod.startswith(target):
                del sys.modules[mod]

# OPTIMIZATION: Order matters. Most common patterns (URLs, emails, filenames) 
# are placed first to speed up regex evaluation by matching earlier.
LOG_EXCLUDED_PATTERN = re.compile(
    # Dates (YYYY-MM-DD, etc.)
    r'\b(?:\d{4}[-/\.]\d{2}[-/\.]\d{2}|\d{2}[-/\.]\d{2}[-/\.]\d{4})\b|'
    # Times (10:30 AM, 14:00, etc.)
    r'\b(?:[01]?\d|2[0-3]):[0-5]\d(?:\s?[aApP][mM])?\b|'
    # IPv4 Addresses
    r'\b(?:\d{1,3}\.){3}\d{1,3}\b|'
    # IPv6 Addresses
    r'(?<![a-zA-Z0-9])(?:[a-fA-F0-9]{0,4}:){2,7}[a-fA-F0-9]{1,4}(?![a-zA-Z0-9])|'
    # MAC Addresses
    r'\b(?:[0-9A-Fa-f]{2}[:-]){5}(?:[0-9A-Fa-f]{2})\b'
)

HIGH_PRIORITY_ENTITIES_PATTERN = re.compile(
    # URLs
    r'https?://[a-zA-Z0-9./_?&=-]+|'                                                                # URLs
    # Emails
    r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}|'                                              # Emails
    # Filenames & Extensions (.pdf, file.mp4)
    r'(?<![a-zA-Z0-9])[a-zA-Z0-9_.-]*\.[a-zA-Z0-9]{2,5}\b|'
    # Phone Numbers
    r'\b(?:\+?\d{1,3}[-.\s]?)?\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}\b|'
    # Social Media Handles / Hashtags
    r'(?<![a-zA-Z0-9])[@#][a-zA-Z0-9_]+\b|'
    # Version Numbers & Kernels
    r'\b(?:[vV]\d+(?:\.\d+)+(?:-[a-zA-Z0-9]+(?:[.-][a-zA-Z0-9]+)*)?|\d+(?:\.\d+)+-[a-zA-Z0-9]+(?:[.-][a-zA-Z0-9]+)*|\d+(?:\.\d+){2,})\b|'
    # Multi-word Proper Nouns / Products
    r'\b[A-Z][a-zA-Z0-9]*(?: [A-Z0-9][a-zA-Z0-9]*){1,3}\b|'
    # UUIDs/GUIDs
    r'\b[0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{12}\b|'
    # Ethereum Addresses
    r'\b0x[a-fA-F0-9]{40}\b|'
    # Bitcoin Addresses
    r'\b(?:[13][a-km-zA-HJ-NP-Z1-9]{25,34}|bc1[a-zA-HJ-NP-Z0-9]{39,59})\b'
)

CODE_ENTITIES_PATTERN = re.compile(
    # snake_case / kebab-case
    r'(?<![a-zA-Z0-9])[a-zA-Z0-9]+(?:[_-][a-zA-Z0-9]+)+(?![a-zA-Z0-9])|'
    # CamelCase / PascalCase
    r'(?<![a-zA-Z0-9])(?:[a-z0-9]+[A-Z][a-zA-Z0-9]*|[A-Z]+[a-z0-9]+[A-Z][a-zA-Z0-9]*|[A-Z]{2,}[a-z0-9][a-zA-Z0-9]*)(?![a-zA-Z0-9])'
)

WORD_PATTERN = re.compile(r'\b[a-z0-9]{3,}\b')

CODE_EXTENSIONS = {".py", ".js", ".jsx", ".ts", ".tsx", ".html", ".css", ".json", ".xml", ".yaml", ".yml", ".c", ".cpp", ".h", ".java", ".cs", ".go", ".rs", ".rb", ".php", ".sh", ".bat", ".ps1", ".sql", ".ini"}
PLAIN_TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".log", ".htm"} | CODE_EXTENSIONS

# Pre-compiled regex patterns used for heuristically filtering out meaningless text tokens.
# Filters 5+ consecutive consonants (likely random hashes)
CONS_PATTERN = re.compile(r'[bcdfghjklmnpqrstvwxz]{5,}')

# Filters alphanumeric hashes with embedded digits (e.g., tD2ar)
SURROUNDED_DIGIT_PATTERN = re.compile(r'[A-Za-z][0-9]+[A-Za-z]')

# Splits strings by delimiters to inspect individual parts
SPLIT_PATTERN = re.compile(r'[\.\-\_\/]')

# Matches semantic version strings (e.g., v1.0, 1.2.3)
SEMANTIC_VERSION_PATTERN = re.compile(r'^[vV]?\d+(?:\.\d+)+$')

# Matches common short filenames (e.g., ui.js, me.md)
SHORT_FILENAME_PATTERN = re.compile(r'^[a-zA-Z0-9]{1,2}\.[a-z]{2,4}$')

# --- Patterns for Metadata & Tagging ---
# Extracts alphanumeric words from path parts for tagging.
TAG_WORD_PATTERN = re.compile(r'[a-zA-Z0-9]+')

# Matches YYYY:MM:DD format in EXIF date strings.
EXIF_DATE_PATTERN = re.compile(r"(\d{4}):(\d{2}):(\d{2})")

# Matches encrypted file extensions like .crypt12
CRYPT_EXT_PATTERN = re.compile(r"^\.crypt\d{2,}$")

def extract_top_keywords(text_data: str, max_words: int = None, is_log: bool = False, extra_entities: list = None) -> str:
    if max_words is None:
        max_words = int(load_config().get("text_extraction_limit", 300))

    def is_meaningful(token: str, is_high_priority: bool = False) -> bool:
        max_len = 150 if is_high_priority else 50
        if len(token) > max_len: 
            return False
        t_lower = token.lower()
        if t_lower in STOP_WORDS: 
            return False
            
        if not is_high_priority and CONS_PATTERN.search(t_lower): 
            return False
            
        parts = [p for p in SPLIT_PATTERN.split(token) if p]
        if not parts: 
            return False # String is just delimiters/symbols
            
        # Skip meaningless minified property sequences (e.g., a.j.ga, n.Pa, c.1g, U.7f)
        if len(parts) > 1 and all(len(p) <= 2 and p.isalnum() for p in parts):
            # Allow semantic versions (e.g., v1.0, 1.2.3)
            if SEMANTIC_VERSION_PATTERN.match(token):
                pass
            # Allow common short lowercase filenames (e.g., ui.js, me.md)
            elif SHORT_FILENAME_PATTERN.match(token):
                pass
            else:
                return False
                
        # Skip isolated 1 or 2 character tokens matched by special patterns (e.g., 8J, 4A, .4Q)
        if len(parts) == 1 and len(token.strip('. \n\r\t')) <= 2:
            return False
            
        is_special = is_high_priority or '@' in token or '://' in token
        for part in parts:
            # Exclude minified random hashes containing digits inside letters (e.g., tD2ar, ddj0n)
            if not is_special and len(part) > 4 and SURROUNDED_DIGIT_PATTERN.search(part):
                return False
                
        return True

    code_entities = CODE_ENTITIES_PATTERN.findall(text_data)
    high_priority_entities = HIGH_PRIORITY_ENTITIES_PATTERN.findall(text_data)
    if not is_log:
        high_priority_entities.extend(LOG_EXCLUDED_PATTERN.findall(text_data))
    if extra_entities:
        high_priority_entities.extend(extra_entities)

    raw_counts = Counter()
    for ent in code_entities:
        if is_meaningful(ent):
            raw_counts[ent] += 1
            
    for ent in high_priority_entities:
        if is_meaningful(ent, is_high_priority=True):
            raw_counts[ent] += 1

    words = WORD_PATTERN.findall(text_data.lower())
    for w in words:
        if is_meaningful(w):
            raw_counts[w] += 1

    # Apply Boosting and Capping logic
    all_counts = Counter()
    
    high_priority_set = set(high_priority_entities)
    unique_high_priority = [ent for ent in raw_counts if ent in high_priority_set]
    unique_high_priority.sort(key=lambda x: raw_counts[x], reverse=True)
    
    strong_id_count = 0
    strong_id_limit = 200
    
    other_id_count = 0
    other_id_limit = 50
    
    proper_noun_count = 0
    proper_noun_limit = 100
    
    ignored_entities = set()
    
    for ent in unique_high_priority:
        # Check if it has no spaces (single-token unique identifier)
        if not any(c.isspace() for c in ent):
            # Is it a strong identifier (email / URL)?
            if '@' in ent or '://' in ent or ent.startswith('www.'):
                if strong_id_count < strong_id_limit:
                    all_counts[ent] = raw_counts[ent] + 100000
                    strong_id_count += 1
                else:
                    ignored_entities.add(ent)
            else:
                if other_id_count < other_id_limit:
                    all_counts[ent] = raw_counts[ent] + 1000
                    other_id_count += 1
                else:
                    ignored_entities.add(ent)
        else:
            if proper_noun_count < proper_noun_limit:
                all_counts[ent] = raw_counts[ent] + 1000
                proper_noun_count += 1
            else:
                ignored_entities.add(ent)
                
    # Add regular words and other code entities (excluding ignored high-priority entities)
    for w, count in raw_counts.items():
        if w not in high_priority_set:
            if w not in ignored_entities:
                all_counts[w] = count

    top_items = [item for item, _ in all_counts.most_common(max_words)]
    return " ".join(top_items)

worker = None

def curate_exemplars_in_memory(faces_list: list, file_dates: dict) -> list:
    """
    Performs 100% in-memory numpy curation on a list of (file_id, embedding) tuples.
    """
    import numpy as np
    if not faces_list:
        return []
    if len(faces_list) <= 15:
        return [emb for _, emb in faces_list]
        
    # Sort faces chronologically based on the file modified date
    sorted_faces = []
    # Sort by date, fallback to empty string
    for fid, emb in sorted(faces_list, key=lambda x: file_dates.get(x[0], "")):
        sorted_faces.append(emb)
        
    embeddings = np.array(sorted_faces, dtype=np.float32)
    norms = np.linalg.norm(embeddings, axis=1, keepdims=True)
    embeddings_norm = embeddings / np.where(norms == 0, 1, norms)

    # A. Centroid (closest to the mean embedding)
    mean_emb = np.mean(embeddings_norm, axis=0)
    mean_norm = np.linalg.norm(mean_emb)
    mean_emb_norm = mean_emb / (mean_norm if mean_norm > 0 else 1.0)
    similarities = np.dot(embeddings_norm, mean_emb_norm)
    centroid_idx = int(np.argmax(similarities))

    # B. Typical representations (top 8 closest to the mean, excluding centroid)
    sorted_by_sim = np.argsort(similarities)[::-1]
    typical_indices = [int(idx) for idx in sorted_by_sim if idx != centroid_idx][:8]

    # C. Diverse boundary faces (up to 6, using Furthest Point Sampling)
    selected_set = set([centroid_idx] + typical_indices)
    boundary_indices = []
    
    median_sim = np.median(similarities)
    candidates = [i for i in range(len(embeddings)) if i not in selected_set and similarities[i] < median_sim]
    if len(candidates) < 6:
        candidates = [i for i in range(len(embeddings)) if i not in selected_set]

    if candidates:
        selected_indices = [centroid_idx] + typical_indices
        for _ in range(min(6, len(candidates))):
            cand_matrix = embeddings_norm[candidates]
            sel_matrix = embeddings_norm[selected_indices]
            sims = np.dot(cand_matrix, sel_matrix.T)
            max_sims = np.max(sims, axis=1)
            best_cand_idx = np.argmin(max_sims)
            best_face_idx = candidates[best_cand_idx]
            boundary_indices.append(int(best_face_idx))
            selected_indices.append(int(best_face_idx))
            candidates.pop(best_cand_idx)

    # D. Timeline-distributed chronological faces (up to 10, evenly spaced)
    indices = np.linspace(0, len(embeddings) - 1, 10, dtype=int)
    timeline_indices = list(dict.fromkeys(indices))

    # Combine in priority order
    all_selected_indices = []
    for idx in [centroid_idx] + typical_indices + boundary_indices + timeline_indices:
        if idx not in all_selected_indices:
            all_selected_indices.append(idx)

    return [sorted_faces[idx] for idx in all_selected_indices[:15]]

def get_or_create_exemplars(person_id: int, conn_or_cursor) -> list:
    """
    Curates and caches up to 15 representative face exemplars for a person profile.
    Uses a fast (under 5ms) in-memory numpy curation to select timeline, centroid, and boundary faces, with a 0ms SQL bypass for small profiles (<= 15 faces).
    """
    from backend.app.state import STATE
    import backend.app.state as app_state
    is_scan_active = app_state.face_scanner_running or app_state.combined_scanner_running
    if (is_scan_active and (STATE.get("face_scanner_stopped") or STATE.get("stopped") or app_state.combined_scanner_stopped)) or STATE.get("cancel_data_operation"):
        return []
        
    from backend.app.utils.cache import EXEMPLAR_CACHE
    cached = EXEMPLAR_CACHE.get(person_id)
    
    cursor = conn_or_cursor if hasattr(conn_or_cursor, "execute") else conn_or_cursor.cursor()
    cursor.execute("SELECT id, file_id, embedding_json FROM faces WHERE person_id = ? AND embedding_json != '[]'", (person_id,))
    known_rows = cursor.fetchall()
    
    if not known_rows:
        return []
        
    current_face_count = len(known_rows)
    if cached and cached.get("count") == current_face_count:
        return cached["embeddings"]
        
    # Bypass logic for small profiles (<= 15 faces)
    if current_face_count <= 15:
        known_embeddings = []
        for face_id, file_id, emb_json in known_rows:
            try:
                emb = json.loads(emb_json)
                if emb and len(emb) == 128:
                    known_embeddings.append(emb)
            except Exception:
                continue
        EXEMPLAR_CACHE.put(person_id, {"count": current_face_count, "embeddings": known_embeddings})
        return known_embeddings

    # For profiles > 15 faces, perform in-memory numpy curation
    log_operation(f"[DEBUG-SCANNER] Curating exemplars in-memory for person {person_id} with {current_face_count} faces...", is_verbose=True)
    faces_list = []
    file_ids = set()
    for face_id, file_id, emb_json in known_rows:
        try:
            emb = json.loads(emb_json)
            if emb and len(emb) == 128:
                faces_list.append((file_id, emb))
                file_ids.add(file_id)
        except Exception:
            continue

    file_ids = list(file_ids)
    file_dates = {}
    with SessionLocal() as s:
        for i in range(0, len(file_ids), 900):
            chunk = file_ids[i:i+900]
            chunk_info = s.query(FileIndex.id, FileIndex.modified).filter(FileIndex.id.in_(chunk)).all()
            for fid, modified in chunk_info:
                file_dates[fid] = str(modified or "")

    curated_embeddings = curate_exemplars_in_memory(faces_list, file_dates)
    EXEMPLAR_CACHE.put(person_id, {"count": current_face_count, "embeddings": curated_embeddings})
    return curated_embeddings


_ocr_engine = None
_ocr_lock = threading.Lock()

def reset_ocr_engine():
    global _ocr_engine
    with _ocr_lock:
        _ocr_engine = None

def get_ocr_engine():
    global _ocr_engine
    with _ocr_lock:
        if _ocr_engine is None:
            try:
                import rapidocr_onnxruntime
                import onnxruntime as ort
                from rapidocr_onnxruntime.utils import UpdateParameters
                from rapidocr_onnxruntime import RapidOCR
            except ImportError as e:
                logging.error(f"Failed to import rapidocr_onnxruntime or onnxruntime: {e}")
                return None

            import os
            import tempfile
            from pathlib import Path
            from backend.app.utils.paths import get_bundled_model_path
            from backend.app.config import load_config
            det_path = get_bundled_model_path("paddleOCR_det.onnx")
            rec_path = get_bundled_model_path("paddleOCR_rec.onnx")
            dict_path = get_bundled_model_path("paddleOCR_dict.txt")
            
            if not Path(det_path).exists() or not Path(rec_path).exists() or not Path(dict_path).exists():
                logging.error(f"OCR model files not found in backend folder. det_path: {det_path}, rec_path: {rec_path}, dict_path: {dict_path}")
                return None

            # Dynamically handle missing config.yaml in packaged environment
            try:
                rapidocr_pkg_dir = Path(os.path.dirname(rapidocr_onnxruntime.__file__))
                default_config = rapidocr_pkg_dir / "config.yaml"
            except Exception:
                rapidocr_pkg_dir = None
                default_config = None

            if not default_config or not default_config.exists():
                temp_dir = Path(tempfile.gettempdir()) / "wabs_rapidocr"
                temp_dir.mkdir(parents=True, exist_ok=True)
                temp_config = temp_dir / "config.yaml"
                
                yaml_content = """Global:
    text_score: 0.5
    use_angle_cls: false
    use_text_det: true
    print_verbose: false
    min_height: 30
    width_height_ratio: 8

Det:
    use_cuda: false
    module_name: ch_ppocr_v3_det
    class_name: TextDetector
    model_path: paddleOCR_det.onnx
    limit_side_len: 736
    limit_type: min
    thresh: 0.3
    box_thresh: 0.5
    max_candidates: 1000
    unclip_ratio: 1.6
    use_dilation: true
    score_mode: fast

Cls:
    use_cuda: false
    module_name: ch_ppocr_v2_cls
    class_name: TextClassifier
    model_path: ch_ppocr_mobile_v2.0_cls_infer.onnx
    cls_image_shape: [3, 48, 192]
    cls_batch_num: 6
    cls_thresh: 0.9
    label_list: ['0', '180']

Rec:
    use_cuda: false
    module_name: ch_ppocr_v3_rec
    class_name: TextRecognizer
    model_path: paddleOCR_rec.onnx
    rec_img_shape: [3, 48, 320]
    rec_batch_num: 6
"""
                try:
                    temp_config.write_text(yaml_content, encoding="utf-8")
                    import rapidocr_onnxruntime.rapid_ocr_api
                    import rapidocr_onnxruntime.utils
                    rapidocr_onnxruntime.rapid_ocr_api.root_dir = temp_dir
                    rapidocr_onnxruntime.utils.root_dir = temp_dir
                    logging.info(f"Recreated missing rapidocr config.yaml at {temp_config} and patched root_dir.")
                except Exception as write_err:
                    logging.error(f"Failed to recreate config.yaml at runtime: {write_err}")
            
            # Load settings
            cfg = load_config()
            cpu_threads = int(cfg.get("ocr_cpu_threads", 4))
            limit_side_len = int(cfg.get("ocr_det_limit_side_len", 736))
            limit_type = str(cfg.get("ocr_det_limit_type", "min"))
            
            # Apply dynamic environment variables as a backup
            if cpu_threads > 0:
                os.environ["OMP_NUM_THREADS"] = str(cpu_threads)
                os.environ["MKL_NUM_THREADS"] = str(cpu_threads)
            
            # Apply dynamic monkeypatch to restrict ONNX Runtime threads and optimize execution providers
            try:
                original_init = ort.InferenceSession.__init__
                def patched_init(self, *args, **kwargs):
                    # Intercept and configure providers
                    providers = kwargs.get('providers')
                    args_list = list(args)
                    providers_in_args = False
                    if providers is None:
                        if len(args_list) > 2:
                            providers = args_list[2]
                            providers_in_args = True
                    
                    try:
                        available = ort.get_available_providers()
                    except Exception:
                        available = []
                        
                    preferred = []
                    if "CUDAExecutionProvider" in available:
                        preferred.append("CUDAExecutionProvider")
                    if "TensorrtExecutionProvider" in available:
                        preferred.append("TensorrtExecutionProvider")
                    if "DmlExecutionProvider" in available:
                        preferred.append("DmlExecutionProvider")
                    if "CPUExecutionProvider" in available:
                        preferred.append("CPUExecutionProvider")
                        
                    if preferred:
                        new_providers = []
                        for p in preferred:
                            if p not in new_providers:
                                new_providers.append(p)
                        if isinstance(providers, (list, tuple)):
                            for p in providers:
                                p_name = p[0] if isinstance(p, tuple) else p
                                if p_name in available and p_name not in new_providers:
                                    new_providers.append(p)
                        if providers_in_args:
                            args_list[2] = new_providers
                            args = tuple(args_list)
                        else:
                            kwargs['providers'] = new_providers

                    sess_opt = kwargs.get('sess_options')
                    if sess_opt is None:
                        if len(args_list) > 1 and isinstance(args_list[1], ort.SessionOptions):
                            sess_opt = args_list[1]
                        else:
                            sess_opt = ort.SessionOptions()
                            if len(args_list) > 1:
                                args_list[1] = sess_opt
                                args = tuple(args_list)
                            else:
                                kwargs['sess_options'] = sess_opt
                    if sess_opt is not None:
                        try:
                            # Disable active spinning in ONNX thread pools to prevent CPU spikes when idle
                            sess_opt.add_session_config_entry("session.intra_op.allow_spinning", "0")
                            sess_opt.add_session_config_entry("session.inter_op.allow_spinning", "0")
                        except Exception:
                            pass
                        if cpu_threads > 0:
                            sess_opt.intra_op_num_threads = cpu_threads
                            sess_opt.inter_op_num_threads = 1
                    return original_init(self, *args, **kwargs)
                ort.InferenceSession.__init__ = patched_init
            except Exception as patch_ex:
                logging.warning(f"Could not monkeypatch onnxruntime InferenceSession.__init__: {patch_ex}")
            
            # Apply monkeypatch to fix RapidOCR's UpdateParameters so we can pass custom keys_path and other settings
            try:
                if hasattr(UpdateParameters, 'update_rec_params'):
                    orig_update_rec = UpdateParameters.update_rec_params
                    def patched_update_rec_params(self, config, rec_dict):
                        if rec_dict:
                            # Normalize rec_keys_path to keys_path
                            if 'rec_keys_path' in rec_dict:
                                rec_dict['keys_path'] = rec_dict.pop('rec_keys_path')
                            # Ensure all keys starting with rec_ have their prefix removed
                            new_rec_dict = {}
                            for k, v in rec_dict.items():
                                if k.startswith('rec_'):
                                    new_rec_dict[k.split('rec_')[1]] = v
                                else:
                                    new_rec_dict[k] = v
                            rec_dict = new_rec_dict
                        return orig_update_rec(self, config, rec_dict)
                    UpdateParameters.update_rec_params = patched_update_rec_params
                else:
                    # Fallback for older versions where update_rec_params is not a standalone helper
                    orig_call = UpdateParameters.__call__
                    def patched_call(self, config, **kwargs):
                        res_config = orig_call(self, config, **kwargs)
                        if 'Rec' in res_config:
                            rec_conf = res_config['Rec']
                            if 'rec_keys_path' in rec_conf:
                                rec_conf['keys_path'] = rec_conf.pop('rec_keys_path')
                            if 'rec_model_path' in rec_conf:
                                rec_conf['model_path'] = rec_conf.pop('rec_model_path')
                        if 'Cls' in res_config:
                            cls_conf = res_config['Cls']
                            if 'cls_model_path' in cls_conf:
                                cls_conf['model_path'] = cls_conf.pop('cls_model_path')
                            if 'cls_label_list' in cls_conf:
                                cls_conf['label_list'] = cls_conf.pop('cls_label_list')
                        return res_config
                    UpdateParameters.__call__ = patched_call
            except Exception as patch_ex:
                logging.warning(f"Could not monkeypatch rapidocr_onnxruntime UpdateParameters: {patch_ex}")

            # Monkeypatch RapidOCR.init_module to support PyInstaller namespace resolving
            try:
                import importlib
                
                def patched_init_module(module_name, class_name):
                    try:
                        module_part = importlib.import_module(module_name)
                    except ImportError:
                        try:
                            module_part = importlib.import_module(f"rapidocr_onnxruntime.{module_name}")
                        except ImportError:
                            raise
                    return getattr(module_part, class_name)
                
                if hasattr(RapidOCR, 'init_module'):
                    RapidOCR.init_module = staticmethod(patched_init_module)
                
                try:
                    import rapidocr_onnxruntime.rapid_ocr_api
                    if hasattr(rapidocr_onnxruntime.rapid_ocr_api, 'init_module'):
                        rapidocr_onnxruntime.rapid_ocr_api.init_module = patched_init_module
                except Exception:
                    pass
            except Exception as patch_ex:
                logging.warning(f"Could not monkeypatch RapidOCR.init_module: {patch_ex}")

            # Check for available GPU providers in ONNX Runtime
            try:
                available = ort.get_available_providers()
            except Exception:
                available = []
            gpu_providers = ["CUDAExecutionProvider", "DmlExecutionProvider", "ROCMExecutionProvider", "TensorrtExecutionProvider"]
            has_gpu = any(p in available for p in gpu_providers)
            
            custom_params = {
                "det_model_path": det_path,
                "rec_model_path": rec_path,
                "rec_keys_path": dict_path,
                "cls_model_path": det_path,  # Bypass for older RapidOCR versions that load classifier unconditionally
                "use_angle_cls": False,
                "det_limit_side_len": limit_side_len,
                "det_limit_type": limit_type
            }
            if has_gpu:
                custom_params["det_use_cuda"] = True
                custom_params["rec_use_cuda"] = True
                logging.info(f"RapidOCR initialized with GPU support (Available: {available})")
            else:
                logging.info(f"RapidOCR initialized on CPU with {cpu_threads} threads (Available: {available})")
                
            try:
                _ocr_engine = RapidOCR(**custom_params)
            except Exception as ex:
                logging.error(f"Failed to instantiate RapidOCR: {ex}")
                _ocr_engine = None
        return _ocr_engine


def _process_unified_scanners(run_index: bool = False, run_face: bool = False, run_object: bool = False, run_document: bool = False):
    try:
        cv2 = _get_cv2()
        fitz = _get_fitz()
        docx = _get_docx()
        pptx = _get_pptx()
        openpyxl = _get_openpyxl()
        filetype = _get_filetype()
        enable_logging = False

        """
        Main routine for processing files through AI scanners and text extraction.
        """
        try:
            from backend.app.config import load_config
            cfg = load_config()
            opencv_threads = int(cfg.get("opencv_cpu_threads", 4))
            if cv2 is not None and opencv_threads > 0:
                cv2.setNumThreads(opencv_threads)
        except Exception:
            pass
            
        if run_index:
            STATE["status"] = "Starting..."
        elif run_document:
            STATE["status"] = "Extracting Text..."
        elif run_face:
            STATE["status"] = "Scanning Faces..."
        elif run_object:
            STATE["status"] = "Classifying Objects..."
        else:
            STATE["status"] = "Starting..."
            
        STATE["stopped"] = False
        if run_face:
            app_state.face_scanner_running = True
            STATE["face_scanner_running"] = True
            app_state.face_scanner_stopped = False
            STATE["face_scanner_stopped"] = False
            STATE["face_scanner_total"] = 0
            STATE["face_scanner_current"] = 0
        if run_object:
            app_state.object_scanner_running = True
            STATE["object_scanner_running"] = True
            app_state.object_scanner_stopped = False
            STATE["object_scanner_stopped"] = False
            STATE["object_scanner_total"] = 0
            STATE["object_scanner_current"] = 0
        if run_document:
            app_state.document_scanner_running = True
            STATE["document_scanner_running"] = True
            STATE["document_scanner_stopped"] = False
            STATE["document_scanner_total"] = 0
            STATE["document_scanner_current"] = 0
        import numpy as np
        from backend.app.utils.paths import get_bundled_model_path, get_ai_db_path
        from backend.app.utils.media import get_cv2_dnn_backends
            
        cfg = load_config()
        enable_logging = cfg.get("enable_logging", False)
        ai_db_path = get_ai_db_path()

        log_operation("Started unified scanners.", user_logs_enabled=enable_logging)
        if enable_logging:
            import logging
            if run_face:
                logging.info("Face scanner process started.")
            if run_object:
                logging.info("Object scanner process started.")
            if run_document:
                logging.info("Document text extraction process started.")

        # --- Object Setup ---
        net, classes, object_threshold = None, None, 0.15
        imagenet_mapping = {}
        if run_object:
            model_path = get_bundled_model_path("mobilenetv2-small.onnx")
            classes_path = get_bundled_model_path("imagenet_classes.txt")
            if Path(model_path).exists() and Path(classes_path).exists():
                net = cv2.dnn.readNetFromONNX(model_path)
                backend_id, target_id = get_cv2_dnn_backends()
                try:
                    net.setPreferableBackend(backend_id)
                    net.setPreferableTarget(target_id)
                    # Test forward pass to verify backend stability
                    test_blob = np.zeros((1, 3, 224, 224), dtype=np.float32)
                    net.setInput(test_blob)
                    net.forward()
                except Exception:
                    # Fallback to CPU
                    try:
                        net.setPreferableBackend(getattr(cv2.dnn, 'DNN_BACKEND_DEFAULT', 0))
                        net.setPreferableTarget(getattr(cv2.dnn, 'DNN_TARGET_CPU', 0))
                    except Exception:
                        pass
            with open(classes_path, 'rt') as f:
                classes = [line.strip() for line in f.readlines()]
            # Load mapping dynamically on-demand
            imagenet_mapping = load_imagenet_mapping()
            object_sensitivity = cfg.get("object_sensitivity", "medium")
            object_threshold = 0.10 if object_sensitivity == "high" else 0.30 if object_sensitivity == "low" else 0.15

        # --- Face Setup ---
        detector, recognizer, clusters, p_count = None, None, {}, 0
        cluster_matrix_norm = None
        cluster_ids_list = []
        new_embs_matrix = None
        new_ids = []
        face_threshold, cluster_threshold = 0.70, 0.55
        # Load YuNet detector if face scanning or object scanning is enabled (for human presence checks)
        if run_face or run_object:
            yunet_path = get_bundled_model_path("face_detection_yunet_2023mar.onnx")
            backend_id, target_id = get_cv2_dnn_backends()

            if Path(yunet_path).exists():
                try:
                    detector = cv2.FaceDetectorYN.create(yunet_path, "", (320, 320), 0.9, 0.3, 5000, backend_id, target_id)
                    # Test detection to verify backend stability
                    test_img = np.zeros((320, 320, 3), dtype=np.uint8)
                    detector.detect(test_img)
                except Exception:
                    backend_id = getattr(cv2.dnn, 'DNN_BACKEND_DEFAULT', 0)
                    target_id = getattr(cv2.dnn, 'DNN_TARGET_CPU', 0)
                    try:
                        detector = cv2.FaceDetectorYN.create(yunet_path, "", (320, 320), 0.9, 0.3, 5000, backend_id, target_id)
                    except Exception:
                        detector = cv2.FaceDetectorYN.create(yunet_path, "", (320, 320))
                
                face_sensitivity = cfg.get("face_sensitivity", "medium")
                face_threshold = 0.55 if face_sensitivity == "high" else 0.85 if face_sensitivity == "low" else 0.70
                detector.setScoreThreshold(face_threshold)

        # Load SFace Recognizer only if face scanning is enabled
        if run_face:
            sface_path = get_bundled_model_path("face_recognition_sface_2021dec.onnx")
            backend_id, target_id = get_cv2_dnn_backends()

            if Path(sface_path).exists():
                try:
                    recognizer = cv2.FaceRecognizerSF.create(sface_path, "", backend_id, target_id)
                except Exception:
                    recognizer = cv2.FaceRecognizerSF.create(sface_path, "")
            else:
                print("Face recognition models not found. Ensure .onnx files are in the backend folder.")
                return

            cluster_sensitivity = cfg.get("face_clustering_sensitivity", "medium")
            cluster_threshold = 0.65 if cluster_sensitivity == "high" else 0.45 if cluster_sensitivity == "low" else 0.55

        # --- DB Initialization ---
        face_processed_ids = set()
        object_processed_ids = set()
        text_processed_ids = set()

        if run_document:
            with SessionLocal() as s:
                s.execute(text("CREATE TABLE IF NOT EXISTS processed_text (file_id INTEGER PRIMARY KEY)"))
                s.execute(text("CREATE VIRTUAL TABLE IF NOT EXISTS file_text_fts USING fts5(file_id UNINDEXED, content)"))
                s.commit()
                text_processed_ids = set(r[0] for r in s.execute(text("SELECT file_id FROM processed_text")).fetchall())

        log_operation("[DEBUG-SCANNER] Initializing SQLite AI database...", is_verbose=True)
        ai_db_path.parent.mkdir(parents=True, exist_ok=True)
        with sqlite3.connect(str(ai_db_path), timeout=15) as conn:
            log_operation("[DEBUG-SCANNER] Setting PRAGMA journal_mode=WAL...", is_verbose=True)
            conn.execute("PRAGMA journal_mode=WAL;")
            cursor = conn.cursor()
            if run_face:
                log_operation("[DEBUG-SCANNER] Creating people table...", is_verbose=True)
                cursor.execute('''CREATE TABLE IF NOT EXISTS people (
                                id INTEGER PRIMARY KEY AUTOINCREMENT,
                                name TEXT DEFAULT 'Unknown Person',
                                thumbnail_file_id INTEGER
                              )''')
                log_operation("[DEBUG-SCANNER] Creating idx_people_name index...", is_verbose=True)
                cursor.execute("CREATE UNIQUE INDEX IF NOT EXISTS idx_people_name ON people(name)")
                log_operation("[DEBUG-SCANNER] Creating faces table...", is_verbose=True)
                cursor.execute('''CREATE TABLE IF NOT EXISTS faces (
                                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                                    person_id INTEGER,
                                    file_id INTEGER,
                                    embedding_json TEXT,
                                    FOREIGN KEY(person_id) REFERENCES people(id)
                                )''')
                log_operation("[DEBUG-SCANNER] Creating idx_faces_person_file index...", is_verbose=True)
                cursor.execute("CREATE INDEX IF NOT EXISTS idx_faces_person_file ON faces(person_id, file_id)")
                log_operation("[DEBUG-SCANNER] Creating idx_faces_unique index...", is_verbose=True)
                cursor.execute("CREATE UNIQUE INDEX IF NOT EXISTS idx_faces_unique ON faces(person_id, file_id, embedding_json)")
                log_operation("[DEBUG-SCANNER] Creating processed_files table...", is_verbose=True)
                cursor.execute('''CREATE TABLE IF NOT EXISTS processed_files (
                                    file_id INTEGER PRIMARY KEY
                                )''')
                log_operation("[DEBUG-SCANNER] Syncing processed_files from faces...", is_verbose=True)
                cursor.execute("INSERT OR IGNORE INTO processed_files (file_id) SELECT DISTINCT file_id FROM faces")
                log_operation("[DEBUG-SCANNER] Querying processed_files...", is_verbose=True)
                cursor.execute("SELECT file_id FROM processed_files")
                face_processed_ids = set(r[0] for r in cursor.fetchall())

                cursor.execute("SELECT COUNT(*) FROM faces")
                if cursor.fetchone()[0] == 0:
                    log_operation("[DEBUG-SCANNER] Database faces table is empty. Purging orphaned face thumbnail cache files...", is_verbose=True)
                    from backend.app.config import get_thumbnail_dir
                    try:
                        thumb_dir = get_thumbnail_dir("faces")
                        if thumb_dir.exists():
                            for f in thumb_dir.glob("person_*.jpg"):
                                try:
                                    f.unlink()
                                except Exception:
                                    pass
                    except Exception as e:
                        log_operation(f"[DEBUG-SCANNER] Error purging face cache: {e}", is_verbose=True)

                log_operation("[DEBUG-SCANNER] Loading all face records in a single query...", is_verbose=True)
                cursor.execute("SELECT person_id, file_id, embedding_json FROM faces WHERE embedding_json != '[]' ORDER BY person_id")
                all_faces_rows = cursor.fetchall()
                
                log_operation(f"[DEBUG-SCANNER] Grouping {len(all_faces_rows)} face records by person_id...", is_verbose=True)
                faces_by_person = {}
                file_ids_needed = set()
                for pid, file_id, emb_json in all_faces_rows:
                    if pid not in faces_by_person:
                        faces_by_person[pid] = []
                    try:
                        emb = json.loads(emb_json)
                        if emb and len(emb) == 128:
                            faces_by_person[pid].append((file_id, emb))
                            file_ids_needed.add(file_id)
                    except Exception:
                        continue
                
                # Fetch dates for those file IDs in bulk
                file_dates = {}
                file_ids_list = list(file_ids_needed)
                log_operation(f"[DEBUG-SCANNER] Fetching dates for {len(file_ids_list)} files...", is_verbose=True)
                with SessionLocal() as s:
                    for i in range(0, len(file_ids_list), 900):
                        chunk = file_ids_list[i:i+900]
                        chunk_info = s.query(FileIndex.id, FileIndex.modified).filter(FileIndex.id.in_(chunk)).all()
                        for fid, modified in chunk_info:
                            file_dates[fid] = str(modified or "")

                log_operation(f"[DEBUG-SCANNER] Curating exemplars for {len(faces_by_person)} people clusters...", is_verbose=True)
                for c_idx, (p_id, faces_list) in enumerate(faces_by_person.items()):
                    if STATE.get("face_scanner_stopped") or STATE.get("stopped") or app_state.combined_scanner_stopped:
                        break
                        
                    if c_idx > 0 and c_idx % 5000 == 0:
                        log_operation(f"[DEBUG-SCANNER] Curated {c_idx} / {len(faces_by_person)} people clusters...", is_verbose=True)
                        
                    clusters[p_id] = curate_exemplars_in_memory(faces_list, file_dates)
                log_operation("[DEBUG-SCANNER] Face embedding preloading and curation completed.", is_verbose=True)

                cursor.execute("SELECT MAX(id) FROM people")
                p_row = cursor.fetchone()
                p_count = p_row[0] if (p_row and p_row[0]) else 0

                # Build initial numpy matrix for clustering
                cluster_embs = []
                for pid, embs in clusters.items():
                    cluster_ids_list.extend([pid] * len(embs))
                    cluster_embs.extend(embs)
                if cluster_embs:
                    cm = np.array(cluster_embs)
                    cn = np.linalg.norm(cm, axis=1, keepdims=True)
                    cluster_matrix_norm = cm / np.where(cn == 0, 1, cn)

            if run_object:
                cursor.execute('''CREATE TABLE IF NOT EXISTS processed_objects (file_id INTEGER PRIMARY KEY)''')
                cursor.execute("SELECT COUNT(*) FROM processed_objects")
                if cursor.fetchone()[0] == 0:
                    with SessionLocal() as s:
                        tagged_photos = s.query(FileIndex.id).filter(FileIndex.category == 'photo', FileIndex.tags.like('%object:%')).all()
                        for (p_id,) in tagged_photos:
                            cursor.execute("INSERT OR IGNORE INTO processed_objects (file_id) VALUES (?)", (p_id,))
                cursor.execute("SELECT file_id FROM processed_objects")
                object_processed_ids = set(r[0] for r in cursor.fetchall())
            conn.commit()

        # --- Build File List ---
        log_operation("[DEBUG-SCANNER] Building file list...", is_verbose=True)
        files_to_process = []
        preloaded_cache = {}
        
        if run_index:
            STATE["status"] = "Discovering files..."
            backup_configs = cfg.get("backup_configs", [])
            roots = [Path(c.get("backup_path", "")) for c in backup_configs if c.get("backup_path")]
            valid_roots = [r for r in roots if r.exists() and r.is_dir()]
            for root_path in valid_roots:
                log_operation(f"Scanning backup path: {root_path}", user_logs_enabled=enable_logging)
                if app_state.combined_scanner_stopped or STATE.get("stopped"):
                    break
                for dirpath, _, filenames in os.walk(str(root_path)):
                    if app_state.combined_scanner_stopped or STATE.get("stopped"):
                        break
                    for f in filenames:
                        if app_state.combined_scanner_stopped or STATE.get("stopped"):
                            break
                        files_to_process.append(os.path.join(dirpath, f))
                        if len(files_to_process) % 1000 == 0:
                            STATE["status"] = f"Discovering files... ({len(files_to_process)} found)"
        else:
            with SessionLocal() as s:
                q = s.query(FileIndex.path, FileIndex.id, FileIndex.category, FileIndex.filename)
                categories = []
                if run_face or run_object:
                    categories.append('photo')
                if run_document:
                    categories.extend(['document', 'ebook', 'code', 'other'])
                    if cfg.get("ocr_enabled", False) and 'photo' not in categories:
                        categories.append('photo')
                if categories:
                    q = q.filter(FileIndex.category.in_(categories))
                for p in q.yield_per(5000):
                    if app_state.combined_scanner_stopped or STATE.get("stopped"):
                        break
                    if not run_index and run_document and STATE.get("document_scanner_stopped"):
                        break
                    if not run_index and run_face and STATE.get("face_scanner_stopped"):
                        break
                    if not run_index and run_object and STATE.get("object_scanner_stopped"):
                        break
                        
                    path_str, item_id, cat, filename = p[0], p[1], p[2], p[3]

                    files_to_process.append(path_str)
                    preloaded_cache[path_str] = {"id": item_id, "size": "", "modified": "", "category": cat, "filename": filename}

        total_files = len(files_to_process)
        log_operation(f"[DEBUG-SCANNER] File list built. Total: {total_files} files.", is_verbose=True)

        if run_index:
            STATE["total"] = total_files
            STATE["current"] = 0
            STATE["indexed"] = 0
            STATE["status"] = "Indexing & Scanning..."
            STATE["running"] = True
        if run_face:
            STATE["face_scanner_total"] = total_files
            STATE["face_scanner_current"] = 0
        if run_object:
            STATE["object_scanner_total"] = total_files
            STATE["object_scanner_current"] = 0
        if run_document:
            STATE["document_scanner_total"] = total_files
            STATE["document_scanner_current"] = 0

        # Give frontend time to register that a scan has officially started with 0 progress
        # This prevents the UI from getting stuck at "Calculating..." if the queue is empty
        # or finishes too quickly.
        if not run_index and not (STATE.get("stopped") or app_state.combined_scanner_stopped or STATE.get("face_scanner_stopped")):
            time.sleep(1.5)
            
        processed_count = 0
        
        with sqlite3.connect(str(ai_db_path), timeout=15) as conn:
            conn.execute("PRAGMA journal_mode=WAL;")
            cursor = conn.cursor()
            with SessionLocal() as session:
                # One-time migration/sync: add "ocr" tag to all photos that have entries in processed_text
                try:
                    query_str = """
                        SELECT id, tags FROM files 
                        WHERE id IN (SELECT file_id FROM processed_text)
                        AND category = 'photo'
                        AND (tags IS NULL OR tags = '' OR (',' || tags || ',') NOT LIKE '%,ocr,%')
                    """
                    rows = session.execute(text(query_str)).fetchall()
                    if rows:
                        log_operation(f"Syncing OCR tags: Found {len(rows)} photos with processed text missing the 'ocr' tag.", is_verbose=True)
                        mappings = []
                        for f_id, existing_tags in rows:
                            tag_list = [t.strip() for t in (existing_tags or "").split(",") if t.strip()]
                            if "ocr" not in tag_list:
                                tag_list.append("ocr")
                                mappings.append({"id": f_id, "tags": ",".join(tag_list)})
                        if mappings:
                            session.bulk_update_mappings(FileIndex, mappings)
                            session.commit()
                            log_operation(f"Successfully tagged {len(mappings)} existing OCR-processed photos with the 'ocr' tag.", is_verbose=True)
                except Exception as sync_e:
                    logging.error(f"Failed to sync existing OCR tags: {sync_e}")

                log_operation("[DEBUG-SCANNER] Starting file indexing & scanning loop...", is_verbose=True)
                # Store lightweight info mapping dynamically in chunks of 1000 to avoid loading millions of rows at startup
                file_cache = {} if run_index else preloaded_cache
                
                offline_roots = set()
                for idx, file_str in enumerate(files_to_process):
                    if shared_state.APP_SHUTTING_DOWN or STATE.get("stopped") or app_state.combined_scanner_stopped:
                        break
                    
                    if run_index and file_str not in file_cache:
                        log_operation(f"[DEBUG-SCANNER] Loading file cache chunk starting at index {idx}...", is_verbose=True)
                        # Clear old cache to limit memory usage, and load the next chunk of 1,000 files
                        file_cache.clear()
                        chunk = files_to_process[idx : idx + 1000]
                        db_items = session.query(FileIndex.id, FileIndex.path, FileIndex.size, FileIndex.modified, FileIndex.category, FileIndex.filename).filter(FileIndex.path.in_(chunk)).all()
                        for r in db_items:
                            file_cache[r.path] = {"id": r.id, "size": r.size, "modified": r.modified, "category": r.category, "filename": r.filename}
                        for p in chunk:
                            if p not in file_cache:
                                file_cache[p] = None

                    if idx % 100 == 0:
                        enable_logging = load_config().get("enable_logging", False)
                    if run_document and STATE.get("document_scanner_stopped") and app_state.document_scanner_running:
                        app_state.document_scanner_running = False
                        STATE["document_scanner_running"] = False
                        STATE["document_scanner_current_file"] = ""
                        run_document = False
                        STATE["document_scanner_stopped"] = False
                        app_state.document_scanner_stopped = False

                    if run_face and STATE.get("face_scanner_stopped") and app_state.face_scanner_running:
                        app_state.face_scanner_running = False
                        STATE["face_scanner_running"] = False
                        STATE["face_scanner_current_file"] = ""
                        run_face = False
                        STATE["face_scanner_stopped"] = False
                        app_state.face_scanner_stopped = False

                    if run_object and STATE.get("object_scanner_stopped") and app_state.object_scanner_running:
                        app_state.object_scanner_running = False
                        STATE["object_scanner_running"] = False
                        STATE["object_scanner_current_file"] = ""
                        run_object = False
                        STATE["object_scanner_stopped"] = False
                        app_state.object_scanner_stopped = False

                    if not run_index and not run_document and not run_face and not run_object:
                        break

                    while STATE.get("paused"):
                        if shared_state.APP_SHUTTING_DOWN:
                            break
                        time.sleep(0.5)
                        
                        if run_document and STATE.get("document_scanner_stopped") and app_state.document_scanner_running:
                            app_state.document_scanner_running = False
                            STATE["document_scanner_running"] = False
                            STATE["document_scanner_current_file"] = ""
                            run_document = False
                            STATE["document_scanner_stopped"] = False
                            app_state.document_scanner_stopped = False

                        if run_face and STATE.get("face_scanner_stopped") and app_state.face_scanner_running:
                            app_state.face_scanner_running = False
                            STATE["face_scanner_running"] = False
                            STATE["face_scanner_current_file"] = ""
                            run_face = False
                            STATE["face_scanner_stopped"] = False
                            app_state.face_scanner_stopped = False

                        if run_object and STATE.get("object_scanner_stopped") and app_state.object_scanner_running:
                            app_state.object_scanner_running = False
                            STATE["object_scanner_running"] = False
                            STATE["object_scanner_current_file"] = ""
                            run_object = False
                            STATE["object_scanner_stopped"] = False
                            app_state.object_scanner_stopped = False

                        if app_state.combined_scanner_stopped or (run_index and STATE.get("stopped")):
                            break
                        if not run_index and not run_document and not run_face and not run_object:
                            break

                    if shared_state.APP_SHUTTING_DOWN:
                        break
                    if app_state.combined_scanner_stopped or (run_index and STATE.get("stopped")):
                        break
                    if not run_index and not run_document and not run_face and not run_object:
                        break
                        
                    # Check if path starts with any known offline root
                    is_offline = False
                    for off_root in offline_roots:
                        if file_str.lower().startswith(off_root):
                            is_offline = True
                            break
                    
                    if is_offline:
                        if run_index:
                            STATE["current"] += 1
                        if run_face and not STATE.get("face_scanner_stopped"):
                            STATE["face_scanner_current"] += 1
                        if run_object and not STATE.get("object_scanner_stopped"):
                            STATE["object_scanner_current"] += 1
                        if run_document and not STATE.get("document_scanner_stopped"):
                            STATE["document_scanner_current"] += 1
                        continue

                    if not run_index:
                        cached_info = file_cache.get(file_str)
                        if cached_info:
                            db_item_id = cached_info["id"]
                            category = cached_info["category"]
                            
                            obj_stopped = STATE.get("object_scanner_stopped", False)
                            doc_stopped = STATE.get("document_scanner_stopped", False)
                            face_stopped = STATE.get("face_scanner_stopped", False)
                            
                            needs_text = run_document and not doc_stopped and db_item_id not in text_processed_ids and (
                                category in ['document', 'ebook', 'code', 'other'] or (category == 'photo' and cfg.get("ocr_enabled", False))
                            )
                            needs_face = run_face and not face_stopped and db_item_id not in face_processed_ids and category == 'photo'
                            needs_object = run_object and not obj_stopped and db_item_id not in object_processed_ids and category == 'photo'
                            
                            if not needs_face and not needs_object and not needs_text:
                                if run_face and not STATE.get("face_scanner_stopped"):
                                    STATE["face_scanner_current"] += 1
                                if run_object and not STATE.get("object_scanner_stopped"):
                                    STATE["object_scanner_current"] += 1
                                if run_document and not STATE.get("document_scanner_stopped"):
                                    STATE["document_scanner_current"] += 1
                                continue

                    file = Path(file_str)
                    if run_index:
                        STATE["current"] += 1
                        STATE["current_file"] = str(file)
                    if run_face and not STATE.get("face_scanner_stopped"):
                        STATE["face_scanner_current"] += 1
                        STATE["face_scanner_current_file"] = file.name
                    if run_object and not STATE.get("object_scanner_stopped"):
                        STATE["object_scanner_current"] += 1
                        STATE["object_scanner_current_file"] = file.name
                    if run_document and not STATE.get("document_scanner_stopped"):
                        STATE["document_scanner_current"] += 1
                        STATE["document_scanner_current_file"] = file.name
                        
                    if not file.exists():
                        try:
                            anchor = file.anchor
                            if anchor and not Path(anchor).exists():
                                offline_roots.add(anchor.lower())
                            elif not file.parent.exists():
                                offline_roots.add(str(file.parent).lower() + os.sep)
                        except Exception:
                            pass
                        continue

                    # --- 1. Indexing Phase ---
                    cached_info = file_cache.get(file_str)
                    db_item_id = cached_info["id"] if cached_info else None
                    db_item = None
                    
                    if run_index:
                        try:
                            f_stat = file.stat()
                            modified_time = time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(f_stat.st_mtime))
                            file_size = str(f_stat.st_size)
                            
                            if not cached_info:
                                category = classify(file.suffix)
                                metadata, extra_tags = extract_metadata_for_file(file, category)
                                tags = build_tags(metadata, category, file.suffix, file)
                                if extra_tags: tags = ",".join(set(tags.split(",") + extra_tags))
                                
                                db_item = FileIndex(
                                    path=str(file), filename=file.name, category=category,
                                    size=file_size, modified=modified_time, extension=file.suffix,
                                    tags=tags, metadata_json=json.dumps(metadata)
                                )
                                session.add(db_item)
                                session.flush()
                                STATE["indexed"] += 1
                                if STATE["indexed"] % 500 == 0:
                                    session.commit()
                                db_item_id = db_item.id
                                file_cache[file_str] = {"id": db_item.id, "size": file_size, "modified": modified_time, "category": category, "filename": file.name}
                                cached_info = file_cache[file_str]
                            else:
                                if cached_info["size"] != file_size or cached_info["modified"] != modified_time:
                                    db_item = session.get(FileIndex, db_item_id)
                                    category = classify(file.suffix)
                                    metadata, extra_tags = extract_metadata_for_file(file, category)
                                    tags = build_tags(metadata, category, file.suffix, file)
                                    if extra_tags: tags = ",".join(set(tags.split(",") + extra_tags))
                                    db_item.size = file_size
                                    db_item.modified = modified_time
                                    db_item.metadata_json = json.dumps(metadata)
                                    db_item.tags = tags
                                    
                                    cached_info["size"] = file_size
                                    cached_info["modified"] = modified_time
                                    cached_info["category"] = category
                                    
                                    STATE["indexed"] += 1
                                    if STATE["indexed"] % 500 == 0:
                                        session.commit()
                        except Exception as e:
                            print(f"Index error on {file.name}: {e}")
                            continue

                    if not cached_info:
                        continue
                    category = cached_info["category"]

                    # --- 2. AI Phase ---
                    obj_stopped = STATE.get("object_scanner_stopped", False)
                    doc_stopped = STATE.get("document_scanner_stopped", False)
                    face_stopped = STATE.get("face_scanner_stopped", False)
                    
                    needs_text = run_document and not doc_stopped and db_item_id not in text_processed_ids and (
                        category in ['document', 'ebook', 'code', 'other'] or (category == 'photo' and cfg.get("ocr_enabled", False))
                    )
                    needs_face = run_face and not face_stopped and db_item_id not in face_processed_ids and category == 'photo'
                    needs_object = run_object and not obj_stopped and db_item_id not in object_processed_ids and category == 'photo'
                    
                    if not needs_face and not needs_object and not needs_text:
                        continue
                        
                    img = None
                    decode_scale = 1.0
                    
                    # Fetch full SQLAlchemy object ONLY if we need to modify tags or it wasn't fetched yet
                    if not db_item and (needs_object or (needs_text and category == 'photo')):
                        db_item = session.get(FileIndex, db_item_id)
                        
                    filename_lower = cached_info["filename"].lower() if cached_info["filename"] else ""
                    if "screenshot" in filename_lower or "meme" in filename_lower:
                        if needs_face: cursor.execute("INSERT OR IGNORE INTO processed_files (file_id) VALUES (?)", (db_item_id,))
                        if needs_object: cursor.execute("INSERT OR IGNORE INTO processed_objects (file_id) VALUES (?)", (db_item_id,))
                        processed_count += 1
                        if processed_count % 500 == 0:
                            conn.commit()
                            session.commit()
                        continue

                    # --- 3. OPTIMIZATION: Read image ONCE from disk for ML models & OCR ---
                    if (needs_face or needs_object or (needs_text and category == 'photo')) and img is None:
                        try:
                            # OPTIMIZATION: Always open and decode the file exactly once to feed all active scanners (Face, Object, OCR) to prevent disk I/O and processing bottlenecks.
                            from PIL import Image, ImageOps
                            start_decode_time = time.perf_counter()
                            width, height = 0, 0
                            orientation = None
                            ocr_enabled = cfg.get("ocr_enabled", False)
                            try:
                                with Image.open(file) as pil_img:
                                    width, height = pil_img.size
                                    exif = pil_img.getexif()
                                    if exif:
                                        orientation = exif.get(274)
                                    
                                    # If the image has transparency (PNG/GIF/WebP etc.), composite it on white immediately.
                                    # This avoids double disk reads and uses PIL's fast composite on white.
                                    if pil_img.mode in ('RGBA', 'LA') or (pil_img.mode == 'P' and 'transparency' in pil_img.info):
                                        bg = Image.new("RGB", pil_img.size, (255, 255, 255))
                                        if pil_img.mode == 'P':
                                            pil_img = pil_img.convert('RGBA')
                                        mask = pil_img.split()[3] if pil_img.mode == 'RGBA' else pil_img.split()[1]
                                        bg.paste(pil_img, mask=mask)
                                        bg = ImageOps.exif_transpose(bg)
                                        img = cv2.cvtColor(np.array(bg), cv2.COLOR_RGB2BGR)
                                        decode_scale = 1.0
                            except Exception:
                                pass
                            
                            # Filter out very small images/icons from face scans only (e.g., width < 100 or height < 100)
                            if width > 0 and height > 0 and (width < 100 or height < 100):
                                if needs_face:
                                    cursor.execute("INSERT OR IGNORE INTO processed_files (file_id) VALUES (?)", (db_item_id,))
                                    needs_face = False
                                
                                if not needs_face and not needs_object and not needs_text:
                                    elapsed_decode = (time.perf_counter() - start_decode_time) * 1000
                                    log_operation(f"[DEBUG-SCANNER] Checked and skipped tiny file: {file.name} (dimensions={width}x{height}) in {elapsed_decode:.3f} ms", is_verbose=True)
                                    processed_count += 1
                                    if processed_count % 500 == 0:
                                        conn.commit()
                                        session.commit()
                                    continue

                            # If not already loaded (i.e. didn't have transparency), proceed with fast OpenCV decode
                            if img is None:
                                img_array = np.fromfile(str(file), np.uint8)
                                ocr_enabled = cfg.get("ocr_enabled", False)
                                
                                # If JPEG and very large, use scale-on-decode flags
                                if width > 0 and height > 0 and file.suffix.lower() in ('.jpg', '.jpeg'):
                                    max_dim = max(width, height)
                                    if ocr_enabled:
                                        # OCR needs higher resolution, only downscale by 1/2 if image is extremely large (>= 3000px)
                                        if max_dim >= 3000:
                                            img = cv2.imdecode(img_array, cv2.IMREAD_REDUCED_COLOR_2)
                                            decode_scale = 0.5
                                    else:
                                        # OCR disabled: aggressive downscaling for face and object classification only
                                        if max_dim >= 3200:
                                            img = cv2.imdecode(img_array, cv2.IMREAD_REDUCED_COLOR_4)
                                            decode_scale = 0.25
                                        elif max_dim >= 1600:
                                            img = cv2.imdecode(img_array, cv2.IMREAD_REDUCED_COLOR_2)
                                            decode_scale = 0.5
                                
                                if img is None:
                                    img = cv2.imdecode(img_array, cv2.IMREAD_COLOR)
                                    
                                # Rotate image using native OpenCV based on EXIF orientation tag if imdecode succeeded
                                if img is not None and orientation in (3, 6, 8):
                                    if orientation == 3:
                                        img = cv2.rotate(img, cv2.ROTATE_180)
                                    elif orientation == 6:
                                        img = cv2.rotate(img, cv2.ROTATE_90_CLOCKWISE)
                                    elif orientation == 8:
                                        img = cv2.rotate(img, cv2.ROTATE_90_COUNTERCLOCKWISE)

                            if img is None:
                                try:
                                    with Image.open(file) as pil_img:
                                        pil_img = ImageOps.exif_transpose(pil_img)
                                        if pil_img.mode in ('RGBA', 'LA') or (pil_img.mode == 'P' and 'transparency' in pil_img.info):
                                            bg = Image.new("RGB", pil_img.size, (255, 255, 255))
                                            if pil_img.mode == 'P':
                                                pil_img = pil_img.convert('RGBA')
                                            mask = pil_img.split()[3] if pil_img.mode == 'RGBA' else pil_img.split()[1]
                                            bg.paste(pil_img, mask=mask)
                                            pil_img = bg
                                        else:
                                            pil_img = pil_img.convert('RGB')
                                        img = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)
                                        decode_scale = 1.0
                                except Exception:
                                    pass
                                    
                            # Check original dimensions again if PIL failed to read them
                            if img is not None and (width == 0 or height == 0):
                                h_dec, w_dec = img.shape[:2]
                                orig_w = int(w_dec / decode_scale)
                                orig_h = int(h_dec / decode_scale)
                                if orig_w < 100 or orig_h < 100:
                                    if needs_face:
                                        cursor.execute("INSERT OR IGNORE INTO processed_files (file_id) VALUES (?)", (db_item_id,))
                                        needs_face = False
                                    
                                    if not needs_face and not needs_object and not needs_text:
                                        img = None
                                        elapsed_decode = (time.perf_counter() - start_decode_time) * 1000
                                        log_operation(f"[DEBUG-SCANNER] Decoded and skipped tiny file: {file.name} (dimensions={orig_w}x{orig_h}) in {elapsed_decode:.3f} ms", is_verbose=True)
                                        processed_count += 1
                                        if processed_count % 500 == 0:
                                            conn.commit()
                                            session.commit()
                                        continue

                            # Log decoding performance
                            if img is not None:
                                elapsed_decode = (time.perf_counter() - start_decode_time) * 1000
                                log_operation(f"[DEBUG-SCANNER] Decoded file: {file.name} (dimensions={width or img.shape[1]}x{height or img.shape[0]}) in {elapsed_decode:.3f} ms", is_verbose=True)
                            else:
                                elapsed_decode = (time.perf_counter() - start_decode_time) * 1000
                                log_operation(f"[DEBUG-SCANNER] Failed to decode file: {file.name} in {elapsed_decode:.3f} ms", is_verbose=True)

                            # Early downscaling strategy when OCR is enabled
                            if img is not None and ocr_enabled:
                                h, w = img.shape[:2]
                                if max(h, w) > 2000:
                                    scale = 2000.0 / max(h, w)
                                    new_w, new_h = int(w * scale), int(h * scale)
                                    img = cv2.resize(img, (new_w, new_h), interpolation=cv2.INTER_AREA)
                                    decode_scale *= scale
                        except Exception as e:
                            print(f"Error reading image {file.name}: {e}")

                        if img is None:
                            if needs_face: cursor.execute("INSERT OR IGNORE INTO processed_files (file_id) VALUES (?)", (db_item_id,))
                            if needs_object: cursor.execute("INSERT OR IGNORE INTO processed_objects (file_id) VALUES (?)", (db_item_id,))
                            if needs_text: session.execute(text("INSERT OR IGNORE INTO processed_text (file_id) VALUES (:f)"), {"f": db_item_id})
                            processed_count += 1
                            if processed_count % 500 == 0:
                                conn.commit()
                                session.commit()
                            continue

                    # --- 4. Run Face Detector (Moved up to detect people before running object classification) ---
                    if needs_face and detector is not None and img is not None:
                        try:
                            log_operation(f"Starting face detection for file: {file.name}", is_verbose=True)
                            dec_h, dec_w, _ = img.shape
                            face_sensitivity = cfg.get("face_sensitivity", "medium")
                            
                            def run_detection_pass(target_dim):
                                scale = 1.0
                                if max(dec_h, dec_w) > target_dim:
                                    scale = target_dim / max(dec_h, dec_w)
                                    new_w, new_h = int(dec_w * scale), int(dec_h * scale)
                                    det_img = cv2.resize(img, (new_w, new_h))
                                    detector.setInputSize((new_w, new_h))
                                else:
                                    det_img = img
                                    detector.setInputSize((dec_w, dec_h))
                                try:
                                    success, faces = detector.detect(det_img)
                                    if success and faces is not None:
                                        scaled_faces = faces.copy()
                                        if scale != 1.0:
                                            scaled_faces[:, :14] /= scale
                                        return scaled_faces
                                except Exception:
                                    pass
                                return None

                            if face_sensitivity == "high":
                                faces_high = run_detection_pass(1024)
                                faces_low = run_detection_pass(320)
                                
                                combined = []
                                if faces_high is not None:
                                    combined.append(faces_high)
                                if faces_low is not None:
                                    combined.append(faces_low)
                                    
                                if combined:
                                    all_faces = np.vstack(combined)
                                    bboxes = [[int(f[0]), int(f[1]), int(f[2]), int(f[3])] for f in all_faces]
                                    scores = [float(f[14]) for f in all_faces]
                                    indices = cv2.dnn.NMSBoxes(bboxes, scores, score_threshold=face_threshold, nms_threshold=0.4)
                                    if len(indices) > 0:
                                        flat_indices = np.array(indices).flatten().tolist()
                                        faces = all_faces[flat_indices]
                                    else:
                                        faces = None
                                else:
                                    faces = None
                            else:
                                faces = run_detection_pass(800)
                                if faces is None:
                                    faces = run_detection_pass(320)

                            num_detected = len(faces) if faces is not None else 0
                            log_operation(f"Face detector found {num_detected} face(s) in {file.name}", is_verbose=True)

                            if faces is not None:
                                for face_idx, face in enumerate(faces):
                                    if shared_state.APP_SHUTTING_DOWN or STATE.get("face_scanner_stopped") or app_state.combined_scanner_stopped or STATE.get("stopped"):
                                        break
                                    face_align = recognizer.alignCrop(img, face)
                                    face_feature = recognizer.feature(face_align)
                                    embedding = face_feature[0].tolist()

                                    best_match_id = None
                                    best_sim = -1.0
                                    
                                    emb_np = np.array(embedding, dtype=np.float32)
                                    emb_norm = np.linalg.norm(emb_np)
                                    emb_np_norm = emb_np / emb_norm if emb_norm > 0 else emb_np
                                    
                                    if cluster_matrix_norm is not None:
                                        similarities = np.dot(cluster_matrix_norm, emb_np_norm)
                                        max_idx = np.argmax(similarities)
                                        max_sim = similarities[max_idx]
                                        
                                        best_sim = float(max_sim)
                                        best_match_id_candidate = cluster_ids_list[max_idx]
                                        log_operation(f"Comparing face #{face_idx} in {file.name} against {len(clusters)} clusters: best similarity {best_sim:.4f} with person_id {best_match_id_candidate}", is_verbose=True)
                                        if max_sim > cluster_threshold:
                                            best_match_id = best_match_id_candidate
                                            
                                    if new_embs_matrix is not None:
                                        new_similarities = np.dot(new_embs_matrix, emb_np_norm)
                                        max_new_idx = np.argmax(new_similarities)
                                        max_new_sim = new_similarities[max_new_idx]
                                        if max_new_sim > cluster_threshold and max_new_sim > best_sim:
                                            best_match_id = new_ids[max_new_idx]
                                            best_sim = float(max_new_sim)
                                            
                                    if best_match_id is None:
                                        while True:
                                            p_count += 1
                                            cursor.execute("INSERT OR IGNORE INTO people (name) VALUES (?)", (f"Unknown Person #{p_count}",))
                                            if cursor.rowcount > 0:
                                                best_match_id = cursor.lastrowid
                                                break
                                        log_operation(f"Face #{face_idx} in {file.name}: similarity below threshold ({best_sim:.4f} vs threshold {cluster_threshold}). Mapped to new person_id: {best_match_id}", is_verbose=True)
                                        clusters[best_match_id] = [embedding]
                                        if new_embs_matrix is None:
                                            new_embs_matrix = np.array([emb_np_norm], dtype=np.float32)
                                        else:
                                            new_embs_matrix = np.vstack([new_embs_matrix, emb_np_norm])
                                        new_ids.append(best_match_id)
                                    else:
                                        log_operation(f"Face #{face_idx} in {file.name}: Matched existing person_id: {best_match_id} (sim: {best_sim:.4f})", is_verbose=True)
                                        if len(clusters[best_match_id]) < 15:
                                            clusters[best_match_id].append(embedding)
                                            if new_embs_matrix is None:
                                                new_embs_matrix = np.array([emb_np_norm], dtype=np.float32)
                                            else:
                                                new_embs_matrix = np.vstack([new_embs_matrix, emb_np_norm])
                                            new_ids.append(best_match_id)
                                    cursor.execute("INSERT OR IGNORE INTO faces (person_id, file_id, embedding_json) VALUES (?, ?, ?)",
                                                    (best_match_id, db_item_id, json.dumps(embedding)))
                                    log_operation(f"Detected and mapped face for file: {file.name} to person_id: {best_match_id}", user_logs_enabled=enable_logging, is_verbose=True)
                        except Exception as e:
                            print(f"Face processing error on {file.name}: {e}")
                            
                        cursor.execute("INSERT OR IGNORE INTO processed_files (file_id) VALUES (?)", (db_item_id,))
                        face_processed_ids.add(db_item_id)

                    # --- 5. Run Object Classifier ---
                    if needs_object and net is not None and img is not None:
                        try:
                            log_operation(f"Starting object classification for file: {file.name}", is_verbose=True)
                            # Letterbox resize to preserve aspect ratio on a white canvas
                            h, w = img.shape[:2]
                            target_size = 224
                            scale = min(target_size / w, target_size / h)
                            new_w, new_h = int(w * scale), int(h * scale)
                            resized = cv2.resize(img, (new_w, new_h), interpolation=cv2.INTER_AREA if scale < 1 else cv2.INTER_CUBIC)
                            
                            o_img = np.full((target_size, target_size, 3), 255, dtype=np.uint8)
                            dx = (target_size - new_w) // 2
                            dy = (target_size - new_h) // 2
                            o_img[dy:dy+new_h, dx:dx+new_w] = resized

                            o_img = cv2.cvtColor(o_img, cv2.COLOR_BGR2RGB)
                            o_img = o_img.astype(np.float32) / 255.0
                            o_img -= np.array([0.485, 0.456, 0.406])
                            o_img /= np.array([0.229, 0.224, 0.225])
                            o_img = o_img.transpose(2, 0, 1)
                            o_img = np.expand_dims(o_img, axis=0)
                            o_img = np.ascontiguousarray(o_img)

                            net.setInput(o_img)
                            preds = net.forward().flatten()
                            exp_preds = np.exp(preds - np.max(preds))
                            probs = exp_preds / np.sum(exp_preds)
                            
                            classIds = np.argsort(probs)[-5:][::-1]
                            raw_preds = [(classes[cid].split(',')[0].strip(), float(probs[cid])) for cid in classIds]
                            log_operation(f"Object classifier top-5 predictions for {file.name}: {raw_preds}", is_verbose=True)
                            
                            new_tags = set()
                            
                            # Check if the photo has any faces (detected in this pass or already in DB)
                            has_face = False
                            if db_item_id is not None:
                                cursor.execute("SELECT 1 FROM faces WHERE file_id = ?", (db_item_id,))
                                if cursor.fetchone() is not None:
                                    has_face = True
                                    
                            # If no face is in DB, and YuNet detector is loaded, run a quick face check
                            if not has_face and detector is not None and img is not None:
                                try:
                                    h_img, w_img = img.shape[:2]
                                    scale = 320.0 / max(h_img, w_img) if max(h_img, w_img) > 320 else 1.0
                                    new_w, new_h = int(w_img * scale), int(h_img * scale)
                                    det_img = cv2.resize(img, (new_w, new_h))
                                    detector.setInputSize((new_w, new_h))
                                    success, faces = detector.detect(det_img)
                                    if success and faces is not None and len(faces) > 0:
                                        has_face = True
                                except Exception:
                                    pass

                            # 1. Aggregate high-level category probabilities if mapping loaded
                            if imagenet_mapping:
                                high_level_probs = {}
                                for cid in range(len(probs)):
                                    prob = float(probs[cid])
                                    # Filter out low-probability subclass noise (e.g., flat noise floor of confused models)
                                    if prob <= 0.015:
                                        continue
                                    spec_raw = classes[cid].split(',')[0].strip()
                                    spec_clean = spec_raw.lower().replace(" ", "_")
                                    mapped_tags = imagenet_mapping.get(spec_clean, [])
                                    for tag in mapped_tags:
                                        high_level_probs[tag] = high_level_probs.get(tag, 0.0) + prob
                                        
                                # Add high-level tags exceeding threshold
                                for tag, prob in high_level_probs.items():
                                    # Apply a strict threshold for animal/pet categories if humans/faces are present in the image
                                    req_threshold = max(object_threshold, 0.45) if (tag in ("animal", "mammal", "pet", "dog", "cat") and has_face) else object_threshold
                                    if prob > req_threshold:
                                        new_tags.add(f"object:{tag}")
                                        
                            # 2. Add specific tags exceeding threshold
                            for cid in classIds:
                                prob = probs[cid]
                                spec_raw = classes[cid].split(',')[0].strip()
                                spec_clean = spec_raw.lower().replace(" ", "_")
                                mapped_tags = imagenet_mapping.get(spec_clean, []) if imagenet_mapping else []
                                
                                is_animal = "animal" in mapped_tags
                                req_threshold = max(object_threshold, 0.45) if (is_animal and has_face) else object_threshold
                                
                                if prob > req_threshold:
                                    new_tags.add(f"object:{spec_clean}")
                                    
                            if new_tags:
                                current_tags_set = parse_tags(db_item.tags)
                                for tag in new_tags:
                                    current_tags_set.add(tag)
                                db_item.tags = ",".join(sorted(current_tags_set))
                                log_operation(f"Classified objects {list(new_tags)} for file: {file.name}", user_logs_enabled=enable_logging, is_verbose=True)
                            else:
                                log_operation(f"No objects detected above threshold for file: {file.name}", is_verbose=True)
                        except Exception as e:
                            print(f"ERROR: Failed to classify {file.name}: {e}")
                            
                        cursor.execute("INSERT OR IGNORE INTO processed_objects (file_id) VALUES (?)", (db_item_id,))
                        object_processed_ids.add(db_item_id)

                    # --- 6. FTS Text Extraction (PDFs, TXT, MD, photos, etc.) ---
                    if needs_text:
                        try:
                            extracted_text = ""
                            extra_entities = []
                            ext = file.suffix.lower()
                            
                            # Get configuration scan depth
                            depth = load_config().get("document_scan_depth", "low")
                            if depth == "high":
                                pdf_limit, docx_limit, pptx_limit, xlsx_sheet_limit, xlsx_row_limit, text_limit = 999999, 999999, 999999, 999, 999999, 10000000
                            elif depth == "medium":
                                pdf_limit, docx_limit, pptx_limit, xlsx_sheet_limit, xlsx_row_limit, text_limit = 150, 1500, 100, 5, 500, 200000
                            else: # low
                                pdf_limit, docx_limit, pptx_limit, xlsx_sheet_limit, xlsx_row_limit, text_limit = 15, 500, 50, 3, 200, 50000

                            if category == 'photo':
                                ocr_enabled = cfg.get("ocr_enabled", False)
                                ocr_only_no_ai_tags = cfg.get("ocr_only_no_ai_tags", True)
                                skip_ocr = False
                                
                                log_operation(f"Evaluating OCR necessity for photo: {file.name} (ocr_enabled={ocr_enabled}, ocr_only_no_ai_tags={ocr_only_no_ai_tags})", is_verbose=True)
                                if ocr_only_no_ai_tags:
                                    cursor.execute("SELECT 1 FROM faces WHERE file_id = ?", (db_item_id,))
                                    has_face = cursor.fetchone() is not None
                                    has_object = False
                                    if db_item and db_item.tags:
                                        if "object:" in db_item.tags:
                                            has_object = True
                                    if has_face or has_object:
                                        skip_ocr = True
                                        log_operation(f"Skipping OCR for photo {file.name} because image already has faces/objects", is_verbose=True)
                                
                                if not skip_ocr and img is not None:
                                    h, w = img.shape[:2]
                                    min_dim = min(h, w)
                                    if min_dim < 20:
                                        skip_ocr = True
                                        log_operation(f"Skipping OCR for photo {file.name} because its dimensions ({w}x{h}) are too small (< 20px)", is_verbose=True)
                                    elif min_dim < 150:
                                        # Scale shorter side to 150px to ensure text clarity
                                        scale_factor = 150.0 / min_dim
                                        new_w = int(w * scale_factor)
                                        new_h = int(h * scale_factor)
                                        img = cv2.resize(img, (new_w, new_h), interpolation=cv2.INTER_CUBIC)
                                        
                                        # Pad with a white border to at least 736px to prevent RapidOCR's internal extreme upscaling
                                        h_new, w_new = img.shape[:2]
                                        if h_new < 736 or w_new < 736:
                                            pad_h = max(0, 736 - h_new)
                                            pad_w = max(0, 736 - w_new)
                                            top = pad_h // 2
                                            bottom = pad_h - top
                                            left = pad_w // 2
                                            right = pad_w - left
                                            img = cv2.copyMakeBorder(img, top, bottom, left, right, cv2.BORDER_CONSTANT, value=[255, 255, 255])
                                            log_operation(f"Rescaled and padded small photo {file.name} to {img.shape[1]}x{img.shape[0]} for optimal OCR", is_verbose=True)
                                
                                if ocr_enabled and not skip_ocr and img is not None:
                                    log_operation(f"Running OCR engine on photo: {file.name}", is_verbose=True)
                                    ocr_engine = get_ocr_engine()
                                    if ocr_engine is not None:
                                        ocr_results, _ = ocr_engine(img)
                                        if ocr_results:
                                            extracted_text = " ".join([res[1] for res in ocr_results if res and len(res) > 1])
                                            log_operation(f"OCR extracted {len(extracted_text)} characters from photo {file.name}: '{extracted_text[:100]}...'", is_verbose=True)
                                            if extracted_text.strip():
                                                if not db_item:
                                                    db_item = session.get(FileIndex, db_item_id)
                                                if db_item:
                                                    existing_tags = db_item.tags or ""
                                                    tag_list = [t.strip() for t in existing_tags.split(",") if t.strip()]
                                                    if "ocr" not in tag_list:
                                                        tag_list.append("ocr")
                                                        db_item.tags = ",".join(tag_list)
                                        else:
                                            log_operation(f"OCR engine returned no text for photo: {file.name}", is_verbose=True)
                            elif ext == '.pdf' and fitz is not None:
                                with fitz.open(str(file)) as doc:
                                    ocr_enabled = cfg.get("ocr_enabled", False)
                                    ocr_max_pages = int(cfg.get("ocr_max_pages", 3))
                                    log_operation(f"Parsing PDF document {file.name} (pages: {len(doc)}, ocr_enabled={ocr_enabled})", is_verbose=True)
                                    for page_num in range(len(doc)):
                                        if shared_state.APP_SHUTTING_DOWN or STATE.get("document_scanner_stopped") or app_state.combined_scanner_stopped or STATE.get("stopped"):
                                            break
                                        if len(extra_entities) > 1000:
                                            break
                                        page = doc[page_num]
                                        page_text = page.get_text()
                                        
                                        if ocr_enabled and len(page_text.strip()) < 15 and page_num < ocr_max_pages:
                                            try:
                                                log_operation(f"PDF page {page_num} of {file.name} has low text ({len(page_text.strip())} chars). Triggering OCR...", is_verbose=True)
                                                pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
                                                img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)
                                                if pix.n == 4:
                                                    page_img = cv2.cvtColor(img_array, cv2.COLOR_RGBA2BGR)
                                                else:
                                                    page_img = cv2.cvtColor(img_array, cv2.COLOR_RGB2BGR)
                                                
                                                if page_num == 0:
                                                    img = page_img # Cache for face/object reuse
                                                    
                                                ocr_engine = get_ocr_engine()
                                                if ocr_engine is not None:
                                                    ocr_results, _ = ocr_engine(page_img)
                                                    if ocr_results:
                                                        ocr_page_text = " ".join([res[1] for res in ocr_results if res and len(res) > 1])
                                                        if ocr_page_text.strip():
                                                            page_text = ocr_page_text
                                                            log_operation(f"OCR extracted text for PDF page {page_num} of {file.name}: {len(page_text)} chars", is_verbose=True)
                                                            if not db_item:
                                                                db_item = session.get(FileIndex, db_item_id)
                                                            if db_item:
                                                                existing_tags = db_item.tags or ""
                                                                tag_list = [t.strip() for t in existing_tags.split(",") if t.strip()]
                                                                if "ocr" not in tag_list:
                                                                    tag_list.append("ocr")
                                                                    db_item.tags = ",".join(tag_list)
                                            except Exception as ocr_err:
                                                if enable_logging:
                                                    logging.error(f"OCR failed on page {page_num} of {file.name}: {ocr_err}")
                                                    
                                        if page_num < pdf_limit:
                                            extracted_text += page_text + " "
                                        else:
                                            ents = HIGH_PRIORITY_ENTITIES_PATTERN.findall(page_text)
                                            ents.extend(LOG_EXCLUDED_PATTERN.findall(page_text))
                                            for ent in ents:
                                                if '@' in ent or '://' in ent or ent.startswith('www.') or ent.startswith('#'):
                                                    extra_entities.append(ent)
                                    log_operation(f"Completed parsing PDF {file.name}. Total extracted text len: {len(extracted_text)} chars", is_verbose=True)
                            elif ext == '.docx' and docx is not None:
                                log_operation(f"Parsing DOCX document {file.name}", is_verbose=True)
                                word_doc = docx.Document(str(file))
                                for i, p in enumerate(word_doc.paragraphs):
                                    if shared_state.APP_SHUTTING_DOWN or STATE.get("document_scanner_stopped") or app_state.combined_scanner_stopped or STATE.get("stopped"):
                                        break
                                    if len(extra_entities) > 1000:
                                        break
                                    p_text = p.text
                                    if not p_text.strip():
                                        continue
                                    if i < docx_limit:
                                        extracted_text += p_text + " "
                                    else:
                                        ents = HIGH_PRIORITY_ENTITIES_PATTERN.findall(p_text)
                                        ents.extend(LOG_EXCLUDED_PATTERN.findall(p_text))
                                        for ent in ents:
                                            if '@' in ent or '://' in ent or ent.startswith('www.') or ent.startswith('#'):
                                                extra_entities.append(ent)
                                log_operation(f"Completed parsing DOCX {file.name}. Total extracted text len: {len(extracted_text)} chars", is_verbose=True)
                            elif ext == '.pptx' and pptx is not None:
                                log_operation(f"Parsing PPTX document {file.name}", is_verbose=True)
                                ppt_doc = pptx.Presentation(str(file))
                                for i, slide in enumerate(ppt_doc.slides):
                                    if shared_state.APP_SHUTTING_DOWN or STATE.get("document_scanner_stopped") or app_state.combined_scanner_stopped or STATE.get("stopped"):
                                        break
                                    if len(extra_entities) > 1000:
                                        break
                                    slide_text = ""
                                    for shape in slide.shapes:
                                        if hasattr(shape, "text") and shape.text.strip():
                                            slide_text += shape.text.strip() + " "
                                    if i < pptx_limit:
                                        extracted_text += slide_text + " "
                                    else:
                                        ents = HIGH_PRIORITY_ENTITIES_PATTERN.findall(slide_text)
                                        ents.extend(LOG_EXCLUDED_PATTERN.findall(slide_text))
                                        for ent in ents:
                                            if '@' in ent or '://' in ent or ent.startswith('www.') or ent.startswith('#'):
                                                extra_entities.append(ent)
                                log_operation(f"Completed parsing PPTX {file.name}. Total extracted text len: {len(extracted_text)} chars", is_verbose=True)
                            elif ext == '.xlsx' and openpyxl is not None:
                                log_operation(f"Parsing XLSX document {file.name}", is_verbose=True)
                                wb = openpyxl.load_workbook(str(file), data_only=True, read_only=True)
                                for s_idx, sheetname in enumerate(wb.sheetnames):
                                    if shared_state.APP_SHUTTING_DOWN or STATE.get("document_scanner_stopped") or app_state.combined_scanner_stopped or STATE.get("stopped"):
                                        break
                                    if len(extra_entities) > 1000:
                                        break
                                    sheet = wb[sheetname]
                                    for r_idx, row in enumerate(sheet.iter_rows(values_only=True)):
                                        if shared_state.APP_SHUTTING_DOWN or STATE.get("document_scanner_stopped") or app_state.combined_scanner_stopped or STATE.get("stopped"):
                                            break
                                        if len(extra_entities) > 1000:
                                            break
                                        row_text = " ".join([str(cell).strip() for cell in row if cell is not None and str(cell).strip()])
                                        if not row_text:
                                            continue
                                        if s_idx < xlsx_sheet_limit and r_idx < xlsx_row_limit:
                                            extracted_text += row_text + " "
                                        else:
                                            ents = HIGH_PRIORITY_ENTITIES_PATTERN.findall(row_text)
                                            ents.extend(LOG_EXCLUDED_PATTERN.findall(row_text))
                                            for ent in ents:
                                                if '@' in ent or '://' in ent or ent.startswith('www.') or ent.startswith('#'):
                                                    extra_entities.append(ent)
                                wb.close()
                                log_operation(f"Completed parsing XLSX {file.name}. Total extracted text len: {len(extracted_text)} chars", is_verbose=True)
                            elif ext in PLAIN_TEXT_EXTENSIONS:
                                log_operation(f"Parsing plain text document {file.name}", is_verbose=True)
                                with open(str(file), 'r', encoding='utf-8', errors='ignore') as f:
                                    extracted_text = f.read(text_limit)
                                    chunk_size = 128 * 1024  # Read in 128KB chunks
                                    while True:
                                        if shared_state.APP_SHUTTING_DOWN or STATE.get("document_scanner_stopped") or app_state.combined_scanner_stopped or STATE.get("stopped"):
                                            break
                                        if len(extra_entities) > 1000:
                                            break
                                        chunk = f.read(chunk_size)
                                        if not chunk:
                                            break
                                        ents = HIGH_PRIORITY_ENTITIES_PATTERN.findall(chunk)
                                        ents.extend(LOG_EXCLUDED_PATTERN.findall(chunk))
                                        for ent in ents:
                                            if '@' in ent or '://' in ent or ent.startswith('www.') or ent.startswith('#'):
                                                extra_entities.append(ent)
                                log_operation(f"Completed parsing plain text {file.name}. Total extracted text len: {len(extracted_text)} chars", is_verbose=True)
                            
                            extracted_text = extracted_text.strip()
                            if extracted_text or extra_entities:
                                limit = int(load_config().get("text_extraction_limit", 300))
                                log_operation(f"Optimizing extracted text for {file.name} (limit={limit}, text_len={len(extracted_text)} chars)", is_verbose=True)
                                if len(extracted_text.split()) > limit:
                                    optimized_text = extract_top_keywords(extracted_text, max_words=limit, is_log=(ext in ['.log']), extra_entities=extra_entities)
                                else:
                                    optimized_text = extract_top_keywords(extracted_text, is_log=(ext in ['.log']), extra_entities=extra_entities)
                                optimized_text = optimized_text.replace('\x00', ' ')
                                log_operation(f"FTS index insertion for {file.name}. Content: '{optimized_text[:100]}...'", is_verbose=True)
                                try:
                                    with session.begin_nested():
                                        session.execute(text("INSERT INTO file_text_fts (file_id, content) VALUES (:f, :c)"), {"f": db_item_id, "c": optimized_text})
                                except Exception as fts_e:
                                    print(f"FTS index error on {file.name}: {fts_e}")
                        except Exception as e:
                            print(f"FTS extraction error on {file.name}: {e}")
                        finally:
                            try:
                                with session.begin_nested():
                                    session.execute(text("INSERT OR IGNORE INTO processed_text (file_id) VALUES (:f)"), {"f": db_item_id})
                            except Exception:
                                pass
                            text_processed_ids.add(db_item_id)

                    processed_count += 1
                    if processed_count % 500 == 0:
                        session.commit()
                        conn.commit()

                session.commit()
                conn.commit()

            # Prevent frontend UI progress bar from instantly disappearing before registering 100% completion
            if not (STATE.get("stopped") or app_state.combined_scanner_stopped):
                time.sleep(1.5)

    except Exception as e:
        print(f"CRITICAL: Unified Worker Error: {e}")
        traceback.print_exc()
        STATE["status"] = f"Error: {e}"
    finally:
        face_processed = STATE.get("face_scanner_current", 0)
        object_processed = STATE.get("object_scanner_current", 0)
        document_processed = STATE.get("document_scanner_current", 0)
        with app_state.scanner_lock:
            is_stopped = STATE.get("stopped", False) or app_state.combined_scanner_stopped
            if run_face and STATE.get("face_scanner_stopped"):
                is_stopped = True
            if run_object and STATE.get("object_scanner_stopped"):
                is_stopped = True
            if run_document and STATE.get("document_scanner_stopped"):
                is_stopped = True

            if run_index:
                STATE["running"] = False
                
            if is_stopped:
                STATE["status"] = "Stopped"
            elif STATE["status"] in ["Starting...", "Extracting Text...", "Scanning Faces...", "Classifying Objects...", "Discovering files...", "Indexing & Scanning..."]:
                STATE["status"] = "Completed"
            
            app_state.face_scanner_running = False
            app_state.object_scanner_running = False
            app_state.document_scanner_running = False
            app_state.combined_scanner_running = False
    
            STATE["face_scanner_running"] = False
            STATE["object_scanner_running"] = False
            STATE["document_scanner_running"] = False
            
            STATE["face_scanner_current_file"] = ""
            STATE["object_scanner_current_file"] = ""
            STATE["document_scanner_current_file"] = ""
            
            STATE["face_scanner_total"] = 0
            STATE["face_scanner_current"] = 0
            STATE["object_scanner_total"] = 0
            STATE["object_scanner_current"] = 0
            STATE["document_scanner_total"] = 0
            STATE["document_scanner_current"] = 0

            # The 'stopped' flags are intentionally NOT reset here.
            # They are reset at the beginning of the next run to avoid race conditions
            # where the frontend misses the 'stopped=True' signal from a manual stop.


        if enable_logging:
            import logging
            status_val = STATE.get("status", "Completed")
            if run_face:
                logging.info(f"Face scanner process ended. Status: {status_val}, processed {face_processed} files.")
            if run_object:
                logging.info(f"Object scanner process ended. Status: {status_val}, processed {object_processed} files.")
            if run_document:
                logging.info(f"Document text extraction process ended. Status: {status_val}, processed {document_processed} files.")

def classify(ext):
    ext = ext.lower()
    if ext in [".jpg",".jpeg",".png",".webp",".gif",".bmp",".tiff",".raw",".svg",".ico",".xcf", ".dng"]:
        return "photo"
    if ext in [".mp4",".mkv",".avi",".mov",".wmv",".flv",".webm",".m4v",".mpg",".mpeg"]:
        return "video"
    if ext in [".mp3",".wav",".flac",".aac",".ogg",".m4a",".wma",".alac"]:
        return "audio"
    if ext in [".pdf",".doc",".docx",".txt",".rtf",".odt",".xls",".xlsx",".ppt",".pptx",".csv",".md",".log"]:
        return "document"
    if ext in [".epub",".mobi",".azw3",".cbz",".cbr",".chm"]:
        return "ebook"
    if ext in CODE_EXTENSIONS:
        return "code"
    if ext in [".ttf",".otf",".woff",".woff2",".eot"]:
        return "font"
    if ext in [".db",".sqlite",".sqlite3",".mdb",".accdb"] or CRYPT_EXT_PATTERN.match(ext):
        return "database"
    if ext in [".zip",".rar",".7z",".tar",".gz",".bz2",".xz"]:
        return "compressed"
    if ext in [".exe",".msi",".apk",".dmg",".deb",".rpm",".appimage"]:
        return "installer"
    if ext in [".bin",".dat",".iso",".img",".vmdk",".vdi",".qcow2",".mpb"]:
        return "binary"
    return "other"

def build_tags(metadata, category, ext, path_obj=None):
    tags = [category or "other", ext.lower().lstrip('.')]
    if isinstance(metadata, dict):
        if "date" in metadata:
            tags.extend(metadata["date"] if isinstance(metadata["date"], list) else [metadata["date"]])
        if "camera" in metadata:
            tags.append(metadata["camera"].lower())
            
    if path_obj:
        for part in path_obj.parts[:-1]:
            words = TAG_WORD_PATTERN.findall(part)
            tags.extend([w.lower() for w in words if len(w) > 2])
            
        words = TAG_WORD_PATTERN.findall(path_obj.stem)
        tags.extend([w.lower() for w in words if len(w) > 2])
        
    return ",".join(set(tags))

def _normalize_metadata_value(value):
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if isinstance(value, bytes):
        try:
            return value.decode('utf-8', errors='ignore')
        except Exception:
            return str(value)
    if isinstance(value, (list, tuple, set)):
        return [_normalize_metadata_value(v) for v in value]
    if isinstance(value, dict):
        return {str(k): _normalize_metadata_value(v) for k, v in value.items()}
    return str(value)

def extract_metadata_for_file(path, category):
    cv2 = _get_cv2()
    fitz = _get_fitz()
    mutagen = _get_mutagen()
    pefile = _get_pefile()
    metadata = {"file_type": category}
    tags = []

    if category == "photo" and Image is not None:
        try:
            with Image.open(path) as img:
                metadata["format"] = img.format
                metadata["mode"] = img.mode
                metadata["size"] = img.size
                if hasattr(img, "getexif"):
                    exif = img.getexif()
                    if exif:
                        if hasattr(exif, 'get_ifd'):
                            try:
                                gps_ifd = exif.get_ifd(0x8825)
                                if gps_ifd and 1 in gps_ifd and 2 in gps_ifd and 3 in gps_ifd and 4 in gps_ifd:
                                    def _dms_to_dec(dms, ref):
                                        try:
                                            d = float(dms[0])
                                            m = float(dms[1])
                                            s = float(dms[2])
                                            dec = d + (m / 60.0) + (s / 3600.0)
                                            return -dec if ref in ['S', 'W'] else dec
                                        except Exception:
                                            return None
                                    lat = _dms_to_dec(gps_ifd[2], gps_ifd[1])
                                    lon = _dms_to_dec(gps_ifd[4], gps_ifd[3])
                                    if lat is not None and lon is not None:
                                        metadata["gps"] = {"latitude": round(lat, 6), "longitude": round(lon, 6)}
                                        tags.append("location")
                            except Exception:
                                pass
                        mapped = {}
                        for tag_id, value in exif.items():
                            tag = ExifTags.TAGS.get(tag_id, tag_id)
                            if tag in ('MakerNote', 'UserComment', 'PrintImageMatching') or (isinstance(value, bytes) and len(value) > 1000):
                                continue
                            mapped[tag] = _normalize_metadata_value(value)
                        metadata["exif"] = mapped
                        date = mapped.get("DateTimeOriginal") or mapped.get("DateTime")
                        if date:
                            date_text = str(date)
                            metadata["date"] = date_text
                            if m := EXIF_DATE_PATTERN.match(date_text):
                                year, month, day = m.groups()
                                tags.append(f"date:{year}")
                                tags.append(f"date:{year}-{month}")
                                tags.append(f"date:{year}-{month}-{day}")
                                tags.append(f"date:{month}/{day}/{year}")
                        if "Model" in mapped:
                            metadata["camera"] = mapped["Model"]
        except Exception:
            metadata["error"] = "Failed to extract image metadata"
    else:
        metadata["file_name"] = path.name
        metadata["file_extension"] = path.suffix.lower()
        
        if category == "document":
            if path.suffix.lower() == ".pdf":
                if fitz is not None:
                    try:
                        doc = fitz.open(str(path))
                        metadata["pages"] = doc.page_count
                        doc.close()
                    except Exception:
                        pass
            elif path.suffix.lower() in [".txt", ".md", ".csv", ".log"]:
                try:
                    if path.stat().st_size < 10 * 1024 * 1024:
                        with open(path, 'rb') as f:
                            lines = f.read().count(b'\n')
                        metadata["lines"] = lines
                        metadata["pages"] = max(1, (lines + 49) // 50)
                except Exception:
                    pass
        elif category == "code":
            try:
                if path.stat().st_size < 10 * 1024 * 1024:
                    with open(path, 'rb') as f:
                        metadata["loc"] = f.read().count(b'\n')
            except Exception:
                pass
        elif category == "video":
            if cv2 is not None:
                try:
                    hw_params = []
                    if hasattr(cv2, 'CAP_PROP_HW_ACCELERATION') and hasattr(cv2, 'VIDEO_ACCELERATION_ANY'):
                        hw_params = [cv2.CAP_PROP_HW_ACCELERATION, cv2.VIDEO_ACCELERATION_ANY]

                    if hw_params:
                        cap = cv2.VideoCapture(str(path), cv2.CAP_FFMPEG, hw_params)
                        if not cap.isOpened():
                            cap = cv2.VideoCapture(str(path), cv2.CAP_ANY, hw_params)
                    else:
                        cap = cv2.VideoCapture(str(path), cv2.CAP_FFMPEG)

                    if not cap.isOpened():
                        cap = cv2.VideoCapture(str(path))
                    if cap.isOpened():
                        width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
                        height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
                        fps = cap.get(cv2.CAP_PROP_FPS)
                        frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                        
                        if width > 0 and height > 0:
                            metadata["resolution"] = f"{width}x{height}"
                        if fps > 0:
                            metadata["fps"] = round(fps, 2)
                            if frame_count > 0:
                                metadata["duration_seconds"] = round(frame_count / fps, 2)
                        cap.release()
                except Exception:
                    pass
                    
            if mutagen is not None:
                try:
                    vid = mutagen.File(str(path))
                    if vid is not None and hasattr(vid, 'tags') and vid.tags:
                        for key, value in vid.tags.items():
                            if not value:
                                continue
                            val_str = str(value[0]) if isinstance(value, list) else str(value)
                            if len(val_str) > 1000:
                                continue
                            if key.lower() in ['date', 'creation_time', '\xa9day', 'year']:
                                metadata['date'] = val_str
                                if len(val_str) >= 4 and val_str[:4].isdigit():
                                    tags.append(f"date:{val_str[:4]}")
                                break
                except Exception:
                    pass
        elif category == "audio" and mutagen is not None:
            try:
                audio = mutagen.File(str(path))
                if audio is not None:
                    if hasattr(audio, 'info') and audio.info is not None:
                        if hasattr(audio.info, 'length') and audio.info.length:
                            metadata['duration_seconds'] = round(audio.info.length, 2)
                        if hasattr(audio.info, 'bitrate') and audio.info.bitrate:
                            metadata['bitrate'] = audio.info.bitrate
                        if hasattr(audio.info, 'sample_rate') and audio.info.sample_rate:
                            metadata['sample_rate'] = audio.info.sample_rate
                    
                    if hasattr(audio, 'tags') and audio.tags:
                        for key, value in audio.tags.items():
                            if not value:
                                continue
                            val_str = str(value[0]) if isinstance(value, list) else str(value)
                            if len(val_str) > 1000:
                                continue
                            key_lower = key.lower()
                            
                            if key_lower in ['title', 'tit2', '\xa9nam']:
                                metadata['title'] = val_str
                            elif key_lower in ['artist', 'tpe1', '\xa9art']:
                                metadata['artist'] = val_str
                                tags.append(f"artist:{val_str.lower()}")
                            elif key_lower in ['album', 'talb', '\xa9alb']:
                                metadata['album'] = val_str
                            elif key_lower in ['genre', 'tcon', '\xa9gen']:
                                metadata['genre'] = val_str
                                tags.append(f"genre:{val_str.lower()}")
                            elif key_lower in ['date', 'tyer', 'tdrc', '\xa9day']:
                                metadata['date'] = val_str
                                tags.append(f"date:{val_str}")
            except Exception:
                pass
        elif category in ["installer", "binary"] and path.suffix.lower() in [".exe", ".dll", ".sys"] and pefile is not None:
            try:
                pe = pefile.PE(str(path), fast_load=True)
                pe.parse_data_directories(directories=[pefile.DIRECTORY_ENTRY['IMAGE_DIRECTORY_ENTRY_RESOURCE']])
                
                metadata["machine"] = "x64" if pe.FILE_HEADER.Machine == 0x8664 else "x86" if pe.FILE_HEADER.Machine == 0x014c else hex(pe.FILE_HEADER.Machine)
                
                if hasattr(pe, 'FileInfo'):
                    for entry in pe.FileInfo:
                        for structure in entry:
                            if hasattr(structure, 'StringTable'):
                                for st_entry in structure.StringTable:
                                    for key, val in st_entry.entries.items():
                                        try:
                                            k = key.decode('utf-8', 'ignore')
                                            v = val.decode('utf-8', 'ignore').strip()
                                            if v and k in ['FileDescription', 'CompanyName', 'FileVersion', 'ProductName']:
                                                metadata[k] = v
                                                if k == 'CompanyName':
                                                    tags.append(f"company:{v.lower().replace(' ', '')}")
                                        except Exception:
                                            pass
            except Exception:
                pass
                
    return metadata, tags

def llm_classify(metadata, ext, cfg):
    if not ext and not metadata:
        return None
        
    provider_url = cfg.get("ai_provider", "").strip()
    api_key = cfg.get("openai_api_key", "").strip()
    
    if not provider_url or provider_url.lower() == "openai":
        api_url = "https://api.openai.com/v1/chat/completions"
        if not api_key:
            return None
    else:
        api_url = provider_url.rstrip("/")
        if not api_url.endswith("/chat/completions"):
            api_url += "/chat/completions"
        if not api_key:
            api_key = "local-dummy-key"
            
    model = cfg.get("ai_model") or "gpt-4o-mini"
    
    prompt = f"Identify the specific file type or category for a file with the extension '{ext}' and the following metadata: {json.dumps(metadata) if metadata else 'None'}. Reply with a short, highly descriptive category (e.g., 'Source Code', '3D Model', 'Audio', 'Configuration', 'Disk Image'). Maximum 2 words. Only reply with the category name, no punctuation."
    
    data = json.dumps({
        "model": model,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.0,
        "max_tokens": 10
    }).encode("utf-8")
    
    req = urllib.request.Request(api_url, data=data, headers={
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}"
    })
    
    if cfg.get("enable_logging"):
        import logging
        logging.info("Requesting classification data from LLM...")
    try:
        with urllib.request.urlopen(req, timeout=45) as response:
            res = json.loads(response.read().decode())
            category = res.get("choices", [{}])[0].get("message", {}).get("content", "")
            if not category:
                category = ""
            category = category.strip()
            category = re.sub(r'[\'"\.]', '', category).title()
            if category and len(category) <= 25:
                return category
    except Exception as e:
        print(f"LLM Classification error for extension {ext}: {e}")
    return None

def background_ai_categorize(cfg):
    if not cfg.get("ai_enabled"):
        return
    session = SessionLocal()
    try:
        items = session.query(FileIndex).filter(FileIndex.category == "other").limit(500).all()
        if not items:
            return
            
        for idx, item in enumerate(items):
            if shared_state.APP_SHUTTING_DOWN:
                break
            if idx % 100 == 0:
                cfg = load_config()
            if STATE.get("stopped"):
                break
            suggested = llm_classify(json.loads(item.metadata_json or "{}"), item.extension or "", cfg)
            if suggested and suggested.lower() != "other":
                item.category = suggested
                tags = set(item.tags.split(",")) if item.tags else set()
                tags.add(suggested.lower())
                item.tags = ",".join(filter(bool, tags))
                session.add(item)
        session.commit()
    finally:
        session.close()

def background_lazy_hasher():
    """
    Background task to calculate SHA-256 hashes for files with identical sizes to find exact duplicates.
    """
    if STATE.get("hasher_running"):
        return
    STATE["hasher_running"] = True
    STATE["hasher_stopped"] = False
    session = None
    try:
        session = SessionLocal()
        cfg = load_config()
        enable_logging = cfg.get("enable_logging", False)
        log_operation("Started background lazy hasher.", user_logs_enabled=enable_logging)
        dup_sizes_query = session.query(FileIndex.size).filter(
            FileIndex.size != '0', 
            FileIndex.size.isnot(None)
        ).group_by(FileIndex.size).having(func.count(FileIndex.id) > 1).all()
        
        dup_sizes = [row[0] for row in dup_sizes_query]
        if not dup_sizes:
            return
            
        files = []
        for i in range(0, len(dup_sizes), 900):
            chunk = dup_sizes[i:i + 900]
            chunk_files = session.query(FileIndex.id, FileIndex.path, FileIndex.metadata_json).filter(
                FileIndex.size.in_(chunk),
                func.json_extract(FileIndex.metadata_json, '$.sha256').is_(None)
            ).all()
            files.extend(chunk_files)
        
        STATE["hasher_total"] = len(files)
        STATE["hasher_current"] = 0
        STATE["hasher_current_file"] = ""

        updates = 0
        mappings = []
        for idx, (item_id, path, metadata_json) in enumerate(files):
            if shared_state.APP_SHUTTING_DOWN:
                break
            if idx % 100 == 0:
                enable_logging = load_config().get("enable_logging", False)
            while STATE.get("paused"):
                if shared_state.APP_SHUTTING_DOWN:
                    break
                time.sleep(0.5)
                if STATE.get("hasher_stopped") or STATE.get("stopped"):
                    break
            
            if shared_state.APP_SHUTTING_DOWN:
                break
            if STATE.get("hasher_stopped") or STATE.get("stopped"):
                break
            STATE["hasher_current"] += 1
            STATE["hasher_current_file"] = Path(path).name
            try:
                meta = json.loads(metadata_json or "{}")
                file_path = Path(path)
                if file_path.exists() and file_path.is_file():
                    log_operation(f"Lazy hasher: calculating SHA-256 for {file_path.name} (size: {file_path.stat().st_size} bytes)", is_verbose=True)
                    hasher = hashlib.sha256()
                    with open(file_path, 'rb') as f:
                        for chunk in iter(lambda: f.read(4096 * 1024), b""):
                            hasher.update(chunk)
                    meta["sha256"] = hasher.hexdigest()
                    log_operation(f"Lazy hasher: SHA-256 for {file_path.name} is {meta['sha256']}", is_verbose=True)
                    mappings.append({"id": item_id, "metadata_json": json.dumps(meta)})
                    updates += 1
                    
                    if updates >= 500:
                        session.bulk_update_mappings(FileIndex, mappings)
                        session.commit()
                        STATE["duplicates_status_changed_at"] = time.time()
                        mappings = []
                        updates = 0
            except Exception as e:
                pass
        if mappings:
            session.bulk_update_mappings(FileIndex, mappings)
            session.commit()
            STATE["duplicates_status_changed_at"] = time.time()
    except Exception as e:
        print(f"Lazy hasher error: {e}")
    finally:
        STATE["hasher_running"] = False
        STATE["hasher_current_file"] = ""
        STATE["duplicates_status_changed_at"] = time.time()
        STATE["hasher_stopped"] = False
        if session:
            session.close()
        if enable_logging:
            import logging
            logging.info(f"Duplicate verification completed. Processed {STATE.get('hasher_current', 0)} files.")

def run():
    """
    Main routine for discovering and indexing files from configured backup paths.
    """
    try:
        cfg = load_config()
        enable_logging = cfg.get("enable_logging", False)
        backup_configs = cfg.get("backup_configs", [])
        if not backup_configs:
            backup_configs = [{"backup_path": cfg.get("backup_path", "")}]

        roots = [Path(c.get("backup_path", "")) for c in backup_configs if c.get("backup_path")]
        valid_roots = [r for r in roots if r.exists() and r.is_dir()]

        log_operation("Started indexer run.", user_logs_enabled=enable_logging)
        with SessionLocal() as session:
            others = session.query(FileIndex).filter(FileIndex.category == 'other').all()
            updated_count = 0
            for item in others:
                new_cat = classify(item.extension or "")
                if new_cat != "other":
                    item.category = new_cat
                    tags = set(item.tags.split(",")) if item.tags else set()
                    tags.add(new_cat)
                    if "other" in tags:
                        tags.remove("other")
                    item.tags = ",".join(filter(bool, tags))
                    updated_count += 1
            if updated_count > 0:
                session.commit()
                print(f"Self-healed {updated_count} files from 'other' to their new categories.")
    except Exception as e:
        print(f"Pre-scan reclassification error: {e}")

    BATCH_SIZE = 500

    last_root_id = ",".join(str(r) for r in valid_roots)

    resume_index = 0
    if STATE.get("last_root") == last_root_id and STATE.get("indexed", 0) > 0:
        resume_index = STATE["indexed"]
    else:
        STATE["current"] = 0
        STATE["indexed"] = 0

    STATE["running"] = True
    STATE["paused"] = False
    STATE["stopped"] = False
    STATE["status"] = "Scanning"
    STATE["last_root"] = last_root_id
    if resume_index == 0:
        STATE["current"] = 0
        STATE["total"] = 0
    STATE["current_file"] = ""

    if not valid_roots:
        STATE["running"] = False
        STATE["status"] = "No valid backup paths configured or found."
        log_operation(STATE["status"], user_logs_enabled=enable_logging)
        return

    try:
        with SessionLocal() as session:
            STATE["status"] = "Discovering files..."
            
            path_mappings = cfg.get("path_mappings")
            if path_mappings and isinstance(path_mappings, dict):
                mapped_count = 0
                for old_prefix, new_prefix in path_mappings.items():
                    items = session.query(FileIndex).filter(FileIndex.path.startswith(old_prefix)).all()
                    for item in items:
                        mapped_path = item.path.replace(old_prefix, new_prefix, 1)
                        mapped_path = mapped_path.replace('\\', os.sep).replace('/', os.sep)
                        item.path = os.path.normpath(mapped_path)
                        mapped_count += 1
                if mapped_count > 0:
                    session.commit()

            global_excluded_str = cfg.get("global_excluded_paths", "")
            global_excluded_list = [p.strip() for p in global_excluded_str.split(",") if p.strip()]
            
            raw_files = []
            for root_path in valid_roots:
                log_operation(f"Indexing backup path: {root_path}", user_logs_enabled=enable_logging)
                matching_config = next((c for c in backup_configs if c.get("backup_path") and Path(c["backup_path"]) == root_path), {})
                excluded_str = matching_config.get("excluded_paths", "")
                backup_excluded_list = [p.strip() for p in excluded_str.split(",") if p.strip()]
                combined_excluded_list = list(set(global_excluded_list + backup_excluded_list))

                for dirpath, dirnames, filenames in os.walk(str(root_path)):
                    if shared_state.APP_SHUTTING_DOWN or STATE.get("stopped"):
                        break
                    if combined_excluded_list:
                        dirnames[:] = [d for d in dirnames if d not in combined_excluded_list]
                    for f in filenames:
                        if shared_state.APP_SHUTTING_DOWN or STATE.get("stopped"):
                            break
                        raw_files.append(os.path.join(dirpath, f))
                        if len(raw_files) % 1000 == 0:
                            STATE["status"] = f"Discovering files... ({len(raw_files)} found)"
            
            raw_files.sort()

            old_total = STATE.get("total", 0)
            STATE["total"] = len(raw_files)
            STATE["status"] = "Indexing"

            start_offset = 0
            is_update_only = STATE.get("update_only")
            
            existing_paths_set = {r[0] for r in session.query(FileIndex.path).all()}
            
            if is_update_only:
                valid_roots_prefixes = [str(r) + os.sep if not str(r).endswith(os.sep) else str(r) for r in valid_roots]
                raw_files_set = set(raw_files)
                paths_to_delete = []
                unmapped_paths = []
                
                for ep in existing_paths_set:
                    if any(ep.startswith(prefix) for prefix in valid_roots_prefixes):
                        if ep not in raw_files_set:
                            paths_to_delete.append(ep)
                    else:
                        unmapped_paths.append(ep)
                                
                if unmapped_paths and not STATE.get("force_reindex"):
                    STATE["status"] = "Unmapped paths detected! Setup Path Mapping or force re-index."
                    STATE["running"] = False
                    return
                            
                if paths_to_delete:
                    if not STATE.get("force_reindex") and not cfg.get("allow_delete_missing"):
                        STATE["status"] = f"{len(paths_to_delete)} missing files. Enable mapping or force re-index to remove."
                        STATE["running"] = False
                        return

                    STATE["status"] = f"Removing {len(paths_to_delete)} missing/excluded files..."
                    for i in range(0, len(paths_to_delete), 500):
                        chunk = paths_to_delete[i:i+500]
                        from sqlalchemy import text
                        ids_to_delete = [r[0] for r in session.query(FileIndex.id).filter(FileIndex.path.in_(chunk)).all()]
                        if ids_to_delete:
                            try:
                                session.execute(text(f"DELETE FROM processed_text WHERE file_id IN ({','.join(map(str, ids_to_delete))})"))
                                session.execute(text(f"DELETE FROM file_text_fts WHERE file_id IN ({','.join(map(str, ids_to_delete))})"))
                            except Exception:
                                pass
                        session.query(FileIndex).filter(FileIndex.path.in_(chunk)).delete(synchronize_session=False)
                    session.commit()

                files_to_process = [f for f in raw_files if f not in existing_paths_set]
                
                processed_count = len(raw_files) - len(files_to_process)
                STATE["indexed"] = processed_count
                STATE["total"] = len(raw_files)
                STATE["current"] = processed_count
                start_offset = processed_count
            elif resume_index > 0 and resume_index < len(raw_files) and old_total == len(raw_files):
                start_offset = resume_index
                files_to_process = raw_files[start_offset:]
                STATE["total"] = len(raw_files)
                STATE["current"] = start_offset
            elif resume_index >= len(raw_files):
                start_offset = len(raw_files)
                files_to_process = []
                STATE["total"] = len(raw_files)
                STATE["current"] = start_offset
            else:
                files_to_process = raw_files

            for idx, file_str in enumerate(files_to_process):
                if shared_state.APP_SHUTTING_DOWN:
                    break
                file = Path(file_str)

                if STATE["stopped"]:
                    STATE["status"] = "Stopped"
                    break

                while STATE["paused"]:
                    if shared_state.APP_SHUTTING_DOWN:
                        break
                    STATE["status"] = "Paused"
                    time.sleep(0.3)
                    if STATE["stopped"]:
                        break

                if shared_state.APP_SHUTTING_DOWN:
                    break
                if STATE["stopped"]:
                    STATE["status"] = "Stopped"
                    break

                real_idx = start_offset + idx
                if idx % 100 == 0:
                    enable_logging = load_config().get("enable_logging", False)
                STATE["current"] = real_idx + 1
                STATE["current_file"] = str(file)
                STATE["status"] = "Indexing"

                try:
                    modified_time = time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(file.stat().st_mtime))
                    file_size = str(file.stat().st_size)
                    category = classify(file.suffix)
                    
                    if category == "other" and filetype is not None:
                        try:
                            kind = filetype.guess(str(file))
                            if kind:
                                mime = kind.mime
                                if mime.startswith('image/'): category = "photo"
                                elif mime.startswith('video/'): category = "video"
                                elif mime.startswith('audio/'): category = "audio"
                                elif mime.startswith('font/'): category = "font"
                                elif 'pdf' in mime: category = "document"
                                elif 'zip' in mime or 'compressed' in mime or 'tar' in mime: category = "compressed"
                                elif 'executable' in mime or 'msdownload' in mime: category = "installer"
                        except Exception:
                            pass

                    metadata, extra_tags = extract_metadata_for_file(file, category)
                    tags = build_tags(metadata, category, file.suffix, file)
                    if extra_tags:
                        tags = ",".join(set(tags.split(",") + extra_tags))

                    exists = None
                    if not is_update_only and str(file) in existing_paths_set:
                        exists = session.query(FileIndex).filter_by(
                            path=str(file)
                        ).first()

                    if exists:
                        if exists.size != file_size or exists.modified != modified_time or exists.metadata_json != json.dumps(metadata) or exists.tags != tags:
                            exists.size = file_size
                            exists.modified = modified_time
                            exists.metadata_json = json.dumps(metadata)
                            exists.tags = tags
                            session.add(exists)
                            log_operation(f"Updated index for file: {file.name}", user_logs_enabled=enable_logging, is_verbose=True)
                    else:
                        session.add(
                            FileIndex(
                                path=str(file),
                                filename=file.name,
                                category=category,
                                size=file_size,
                                modified=modified_time,
                                extension=file.suffix,
                                tags=tags,
                                metadata_json=json.dumps(metadata)
                            )
                        )
                        log_operation(f"Indexed new file: {file.name}", user_logs_enabled=enable_logging, is_verbose=True)

                    if idx > 0 and idx % BATCH_SIZE == 0:
                        session.commit()
                except Exception as exc:
                    print(f"Indexer error processing {file}: {exc}")
                    traceback.print_exc()
                    session.rollback()
                    continue

            try:
                session.commit()
            except Exception as exc:
                print(f"Indexer final commit failed: {exc}")
                traceback.print_exc()
                session.rollback()

        if cfg.get("ai_enabled") and not STATE.get("stopped"):
            categorization_thread = Thread(target=background_ai_categorize, args=(cfg,), daemon=True)
            categorization_thread.start()
    except Exception as exc:
        try:
            if load_config().get("enable_logging"):
                import logging
                logging.error(f"Indexer error: {exc}")
        except Exception:
            pass
        STATE["status"] = f"Error: {exc}"
    finally:
        STATE["running"] = False
        if STATE["stopped"]:
            STATE["status"] = "Stopped"
        elif STATE["status"] in ["Starting...", "Scanning", "Discovering files...", "Indexing", "Indexing & Scanning..."]:
            STATE["status"] = "Completed"
        STATE["indexed"] = STATE.get("current", 0)
        if enable_logging:
            import logging
            logging.info(f"Indexer run finished. Status: {STATE['status']}, processed {STATE['current']} files.")

def start_indexing():
    global worker

    if STATE["running"]:
        return

    STATE["running"] = True
    STATE["status"] = "Starting..."
    STATE["stopped"] = False
    worker = Thread(target=run, daemon=True)
    worker.start()